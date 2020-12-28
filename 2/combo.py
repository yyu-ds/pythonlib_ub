# -*- coding: utf-8 -*-
"""
Created on Wed May  2 14:37:14 2018

@author: ub71894 (4e8e6d0b), CSG
"""


import pyautogui,time

pl_names=[
'Primus Telecommunications Group, Incorporated',
'Citadel Broadcasting Corporation',
'Motor Coach Industries International, Inc.',
'Beazer Homes USA, Inc.',
'True Temper Sports, Inc.',
'Visteon Corporation',
'CMP Susquehanna Corp',
'BI-LO, LLC (Old)',
'Pliant Corporation (Old)',
'Buffets, Inc. (Old)',
'Smurfit-Stone Container Enterprises (Old)',
'Workflow Management, Inc.',
'Spectrum Brands, Inc.',
'Source Interlink Companies Inc.',
'Atrium Corporation',
'Chemtura Corporation',
'Journal Register Company',
'Newark Group, Inc. (The)',
'Euramax International, Inc.',
'Fleetwood Enterprises Inc.',
'NextMedia Operating, Inc. (Old)',
'Portola Packaging, Inc.',
'Nexstar Broadcasting, Inc.',
'VeraSun Energy Corporation',
'NV Broadcasting, LLC',
'American Achievement Corporation',
'PRC, LLC',
'United Subcontractors Inc.',
'Pacific Lumber Co.',
'Pac-West Telecomm, Inc.',
'Coinmach Service Corp.',
'BearingPoint, Inc.',
'Express Energy Services Operating, LP',
'Hines Horticulture, Inc.',
'SuperMedia Inc.',
'Penton Business Media Holdings, Inc.',
'Emmis Communications Corporation',
'Foamex L.P.',
'iHeartCommunications, Inc.',
'Duane Reade, Inc.',
'Masonite Corporation',
'Six Flags Inc. (Old)',
'Six Flags Inc. (Old)',
'Port Townsend Paper Corporation',
'RBS Global, Inc.',
'Lazy Days R.V. Center, Inc.',
'Axiall Corporation',
'Tribune Media Company',
'Caesars Entertainment Corporation',
'MagnaChip Semiconductor Corporation',
'Young Broadcasting Inc.',
'Appvion, Inc.',
'Recycled Paper Greetings, Inc.',
'Caraustar Industries, Inc.',
'Wellman, Inc.',
'Haights Cross Communications, Inc.',
'Big West Oil, LLC (Old)',
'Aventine Renewable Energy Holdings (Old)',
'ION Media Networks, Inc.',
'Mark IV Industries, Inc. (Reorganized)',
'Key Plastics LLC',
'American Media Operations, Inc.',
'Plastech Engineered Products, Inc.',
'Vertis, Inc. (Old)',
'Suncom Wireless, Inc',
'RathGibson, Inc.',
'Pilgrims Pride Corporation (Old)',
'Citation Corporation',
'Metaldyne Corporation',
'JHT Holdings, Inc.',
'Merisant Worldwide, Inc.',
'Kimball Hill, Inc.',
'Finlay Fine Jewelry Corporation',
'Finlay Fine Jewelry Corporation',
'U.S. Shipping Partners LP',
'Noranda Aluminum Holding Corporation',
'OSI Restaurant Partners, LLC',
'Realogy Group LLC',
'Pierre Foods, Inc. (Old)',
'Headwaters Incorporated',
'Dayton Superior Corporation',
'YRC Worldwide Inc.',
'Regent Broadcasting LLC',
'Hawaiian Telcom Communications, Inc.',
'Electrical Components Intl. Inc. (Old)',
'Lenox Group, Inc.',
'Alion Science and Technology Corp',
'Propex Inc.',
'Movie Gallery, Inc.',
'Movie Gallery, Inc.',
'Builders FirstSource, Inc.',
'NTK Holdings, Inc.',
'Ford Motor Company',
'Jacuzzi Brands Corp.',
'Linens N Things, Inc.',
'Rhodes Companies, LLC (The)',
'Aleris International Inc.',
'Spheris Inc.',
'Momentive Performance Materials Inc. (Old)',
'Simmons Company',
'Quality Home Brands Holdings LLC (OLD)',
'Cygnus Business Media, Inc.',
'Charter Communications Inc.',
'FairPoint Communications, Inc.',
'Lyondell Chemical Company',
'Libbey Glass Inc.',
'Milacron Inc.',
'Chesapeake Corporation',
'Black Gaming, LLC',
'TLC Vision Corporation',
'Freescale Semiconductor, Inc.',
'Sensata Technologies B.V.',
'Affinity Group Holding, Inc.',
'Bally Total Fitness Holding Corporation',
'Bally Total Fitness Holding Corporation',
'Accuride Corporation (Old)',
'U.S. Concrete, Inc.',
'Chem Rx Corporation',
'Uno Restaurant Holdings Corporation',
'EnviroSolutions Holdings, Inc.',
'InSight Health Services Corp.',
'Constar International, Inc.',
'Legends Gaming, LLC (OLD)',
'Neff Corp.',
'Neff Corp.',
'Spansion, LLC',
'Penn Traffic Company',
'Wornick Company (The)',
'Century Aluminum Company',
'Neenah Foundry Company',
'Frontier Airlines, Inc.',
'Quantum Corporation',
'Champion Enterprises, Inc.',
'Readers Digest Association, Inc. ',
'Lear Corporation',
'Stallion Oilfield Holdings, Inc.',
'Cooper-Standard Automotive Inc.',
'Xerium Technologies, Inc.',
'R.H. Donnelley Inc. (Old)']

#%%


time.sleep(10)
for name in pl_names:
    
	pyautogui.moveTo(398, 172)
	pyautogui.click(398, 172)
	pyautogui.typewrite(name)
	pyautogui.press('enter')	

	time.sleep(4)
	pyautogui.click(627, 416)
	time.sleep(3)	

	pyautogui.moveTo(533, 279)
	pyautogui.dragTo(601, 279, button='left')
	pyautogui.hotkey('ctrl', 'c')	
	
	pyautogui.hotkey('alt', 'tab')
	pyautogui.hotkey('ctrl', 'v')
	pyautogui.press('enter')
	pyautogui.hotkey('alt', 'tab')


#%%

pyautogui.position()
asdasd
# -*- coding: utf-8 -*-
"""
Created on Wed Jun 17 17:06:38 2015

@author: ub69401
"""

import re
import os,csv, sqlite3

#define path name to put the database file
dir_name=r'C:\Users\ub69401\Documents\compustat\database'
#define result file name
base_filename='compustat_new'
filename_suffix = '.db'
filename = os.path.join(dir_name, base_filename + filename_suffix)
#create and connect to db file
sqlite_file = filename
con = sqlite3.connect(filename)
cur = con.cursor()

path = r'C:\Users\ub69401\Documents\compustat'
for root, dirs, files in os.walk(path): #go through all files, subdirectories
    for name in files: #loop over files
        if name.endswith((".txt")): #but only open .txt files
            searchObj = re.search( r'_(.*?).txt$' , str(name), re.M|re.I) #strip out the group name
            print(searchObj.group(1)) #check if it is correct
            table_name = searchObj.group(1) #group name becomes the table name
            with open(os.path.join(root,name),"rt",encoding="utf8") as fin: #open each txt for reading
                dr = csv.reader(fin,delimiter='|') #use csv.reader to read file
                columns1 = next(dr) #read in the header line 
                columns2 = columns1[4:] #read header line, ignore first 4 elements as meta-data
                ncol     = len(columns1)-4 #adjust number of columns
                column_list = "(" #build column names and question string
                questionstring ="("
                for n in range(ncol):
                    if n+1 < (ncol):
                        column_list = column_list+str(columns2[n])+", "
                        questionstring = questionstring +"?," 
                    else:
                        column_list = column_list+str(columns2[n])+")"
                        questionstring = questionstring + "?)"
                print(name," ",ncol)
                print(column_list)
                #fin.seek(0) #go back to the beginning of the file
                executestring = "DROP TABLE IF EXISTS "+str(table_name)+";CREATE TABLE "+str(table_name)+" "+column_list+";"
                cur.executescript(executestring) #create table if it does not already exist
                for row in dr:  #write remaining rows into the table
                    to_db = [str(row[n]) for n in range(ncol)]
                    writestring = "INSERT INTO "+str(table_name)+" "+column_list+" VALUES "+questionstring+";"
                    cur.executemany(writestring, (to_db,))
                    con.commit()
con.close()

# -*- coding: utf-8 -*-
"""
Created on Mon Nov  9 16:56:38 2015

@author: ub71894  modified based on ub69401
"""

import re
import os,csv, sqlite3

#define path name to put the database file
dir_name=r'C:\Users\ub71894\Documents\Python Scripts'
#define result file name
base_filename='new2'
filename_suffix = '.db'
filename = os.path.join(dir_name, base_filename + filename_suffix)
#create and connect to db file
sqlite_file = filename
con = sqlite3.connect(filename)
cur = con.cursor()

path = r'C:\Users\ub71894\Documents\Python Scripts'
for root, dirs, files in os.walk(path): #go through all files, subdirectories
    for name in files: #loop over files
        if name.endswith((".txt")): #but only open .txt files
            searchObj = re.search( r'_(.*?).txt$' , str(name), re.M|re.I) #strip out the group name
            print(searchObj.group(1)) #check if it is correct
            table_name = searchObj.group(1) #group name becomes the table name
            with open(os.path.join(root,name),"rt",encoding="utf8") as fin: #open each txt for reading
                dr = csv.reader(fin,delimiter='|') #use csv.reader to read file
                columns1 = next(dr) #read in the header line 
                columns2 = columns1[4:] #read header line, ignore first 4 elements as meta-data
                ncol     = len(columns1)-4 #adjust number of columns
                # column_list = "(" #build column names and question string
                # questionstring ="("
                # for n in range(ncol):
                #     if n+1 < (ncol):
                #         column_list = column_list+str(columns2[n])+", "
                #         questionstring = questionstring +"?," 
                #     else:
                #         column_list = column_list+str(columns2[n])+")"
                #         questionstring = questionstring + "?)"
                # print(name," ",ncol)
                # print(column_list)
                #fin.seek(0) #go back to the beginning of the file
                executestring = "DROP TABLE IF EXISTS "+str(table_name)+";CREATE TABLE "+str(table_name)+" ("+ ','.join(map(str,columns2))+ ");" # modified
                cur.executescript(executestring) #create table if it does not already exist
                for row in dr:  #write remaining rows into the table
                    to_db = [str(row[n]) for n in range(ncol)]
                    writestring = "INSERT INTO "+str(table_name)+" ("+ ','.join(map(str,columns2))+ ") VALUES "+ " ("+ ','.join('?'*len(columns2))+ ");"# modified
                    cur.executemany(writestring, (to_db,))
                    con.commit()
con.close()

# -*- coding: utf-8 -*-
"""
Created on Mon Mar 30 13:43:39 2020

@author: mua2a76
"""

#
# This script has the details and instructions used to build SQLite database and upload Compustat data
# 
#%%
# Step 1: Install SQLite and an editor such as SQLiteStudio
# The below will create a database named insurance032020 and 2 tables to hold the financials for all companies.
# You can check in the editor that the database appears in the file explorer.

#%%

import os, sqlite3, csv
from io import StringIO


dir_name=r'C:\Users\ru87118\Desktop\Compustat\Data'
os.chdir(dir_name)

base_filename='CnI091920_2'
filename_suffix = '.db'
filename = os.path.join(dir_name, base_filename + filename_suffix)
con = sqlite3.connect(filename)


with open(r'C:\Users\ru87118\Desktop\Compustat\Data\SQL\Financials_M_Z_create table.sql', 'r') as file:
    sql_create_Financials_M_Z_table = file.read().replace('\n', ' ')
    
with open(r'C:\Users\ru87118\Desktop\Compustat\Data\SQL\Financials_A_L_create table.sql', 'r') as file:
    sql_create_Financials_A_L_table = file.read().replace('\n', ' ')
    
with open(r'C:\Users\ru87118\Desktop\Compustat\Data\SQL\Company_create table.sql', 'r') as file:
    sql_create_Company_table = file.read()
    
c = con.cursor()

c.execute(sql_create_Financials_M_Z_table)
c.execute(sql_create_Financials_A_L_table)
c.execute(sql_create_Company_table)

con.commit()
c.close()
con.close()

#%%
# Step 2: Import the data in the text files to the two financials tables
#         First, extract the files f_co_afnd1V2 and f_co_afndv2V2 in folders with the same name to avoid special characters in windows folder names

data_dir_A_L = r'C:\Users\ru87118\Desktop\Compustat\Data'
filename_A_L = r'f_co_afnd1V2.txt'
remove_string_A_L = 'H|20200919-07:42:53|co_afnd1V2|6|'

file_A_L = open(os.path.join(data_dir_A_L, filename_A_L),'r').read().replace(remove_string_A_L,'').replace('|',',')
s = StringIO(file_A_L)
# Save a csv for direct manual inspection
with open('afnd1V2.csv', 'w') as f:
    for line in s:
        f.write(line)


table_name = 'Financials_A_L'
con = sqlite3.connect(filename)
cur = con.cursor()

with open(r'C:\Users\ru87118\Desktop\Compustat\Data\afnd1V2.csv',"rt",encoding="utf8") as fin: 
    dr = csv.reader(fin,delimiter=',',quoting=csv.QUOTE_NONE,lineterminator='\n') #use csv.reader to read file
    columns1 = next(dr) #read in the header line 
    ncol     = len(columns1)
    column_list = "(" #build column names and question string
    questionstring ="("
    for n in range(ncol):
        if n+1 < (ncol):
            column_list = column_list+str(columns1[n])+", "
            questionstring = questionstring +"?," 
        else:
            column_list = column_list+str(columns1[n])+")"
            questionstring = questionstring + "?)"
    print(column_list)
    for row in dr:  #write remaining rows into the table
        to_db = [str(row[n]) for n in range(ncol)]         
        writestring = "INSERT INTO "+str(table_name)+" "+column_list+" VALUES "+questionstring+";"
        cur.executemany(writestring, (to_db,))
        con.commit()
con.close()


#%%
# Step 2: Continuation
#

data_dir_M_Z = r'C:\Users\ru87118\Desktop\Compustat\Data'
filename_M_Z = r'f_co_afnd2V2.txt'
remove_string_M_Z = 'H|20200919-07:41:52|co_afnd2V2|6|'

file_M_Z = open(os.path.join(data_dir_M_Z, filename_M_Z),'r').read().replace(remove_string_M_Z,'').replace('|',',')
s = StringIO(file_M_Z)
# Save a csv for direct manual inspection
with open('afnd2V2.csv', 'w') as f:
    for line in s:
        f.write(line)

table_name = 'Financials_M_Z'
con = sqlite3.connect(filename)
cur = con.cursor()

with open(r'C:\Users\ru87118\Desktop\Compustat\Data\afnd2V2.csv',"rt",encoding="utf8") as fin: 
    dr = csv.reader(fin,delimiter=',',quoting=csv.QUOTE_NONE,lineterminator='\n') #use csv.reader to read file
    columns1 = next(dr) #read in the header line 
    ncol     = len(columns1)
    column_list = "(" #build column names and question string
    questionstring ="("
    for n in range(ncol):
        if n+1 < (ncol):
            column_list = column_list+str(columns1[n])+", "
            questionstring = questionstring +"?," 
        else:
            column_list = column_list+str(columns1[n])+")"
            questionstring = questionstring + "?)"
    print(column_list)
    for row in dr:  #write remaining rows into the table
        to_db = [str(row[n]) for n in range(ncol)]         
        writestring = "INSERT INTO "+str(table_name)+" "+column_list+" VALUES "+questionstring+";"
        cur.executemany(writestring, (to_db,))
        con.commit()
con.close()

#%%
# Combine the two tables Financilas_A_L and Fiancials_M_Z into Financials table
# Note that a full join is required here. Sqlite does not support it. This could lead to a minimal loss of data.
con = sqlite3.connect(filename)

with open(r'C:\Users\ru87118\Desktop\Compustat\Data\SQL\Financials_create table.sql', 'r') as file:
    sql_create_Financials_table = file.read().replace('\n', ' ')
    
c = con.cursor()

c.execute(sql_create_Financials_table)

con.commit()
c.close()
con.close()

#%%
# Populate the Company table
data_dir_company = r'C:\Users\ru87118\Desktop\Compustat\Data'
filename_company = r'f_company.txt'
remove_string_company = 'H|20200919-07:32:36|company|1|'

file_company = open(os.path.join(data_dir_company, filename_company),'r').read().replace(remove_string_company,'')
s = StringIO(file_company)
# Save a csv for direct manual inspection
with open('company.txt', 'w') as f:
    for line in s:
        f.write(line)

table_name = 'Company'
con = sqlite3.connect(filename)
cur = con.cursor()

with open(r'C:\Users\ru87118\Desktop\Compustat\Data\company.txt',"rt",encoding="utf8") as fin: 
    dr = csv.reader(fin,delimiter='|',quoting=csv.QUOTE_NONE,lineterminator='\n') #use csv.reader to read file
    columns1 = next(dr) #read in the header line 
    ncol     = len(columns1)
    column_list = "(" #build column names and question string
    questionstring ="("
    for n in range(ncol):
        if n+1 < (ncol):
            column_list = column_list+str(columns1[n])+", "
            questionstring = questionstring +"?," 
        else:
            column_list = column_list+str(columns1[n])+")"
            questionstring = questionstring + "?)"
    print(column_list)
    for row in dr:  #write remaining rows into the table
        to_db = [str(row[n]) for n in range(ncol)]         
        writestring = "INSERT INTO "+str(table_name)+" "+column_list+" VALUES "+questionstring+";"
        cur.executemany(writestring, (to_db,))
        con.commit()
con.close()
#%%
# Construct insurance financials table. This will hold financial data that will be used as input to the python scripts
#
con = sqlite3.connect(filename)

with open(r'C:\Users\ru87118\Desktop\Compustat\Data\SQL\CnIFinancials_create table.sql', 'r') as file:
    sql_create_CnIFinancials_table = file.read().replace('\n', ' ')
    
c = con.cursor()

c.execute(sql_create_CnIFinancials_table)

con.commit()
c.close()
con.close()
#
#
# End of Financial data section
#
#
#%%
# Rating data and related tables
#
data_dir_ratings = r'C:\Users\ru87118\Desktop\Compustat\Data'
filename_ratings = r'spRatingData.txt'

ratingsTableColumns=['ratingDetailId'	,'entitySymbolValue','instrumentSymbolValue'	,'securitySymbolValue',
                     'objectTypeId'	,'orgDebtTypeCode','ratingTypeCode','currentRatingSymbol','ratingSymbol',
                     'ratingDate','creditwatch','outlook','creditwatchDate','outlookDate','priorRatingSymbol',	
                     'priorCreditwatch','priorOutlook','ratingQualifier','regulatoryIndicator','regulatoryQualifier',
                     'ratingActionWord','CWOLActionWord','cwolInd','maturityDate','CUSIP','CINS','ISIN']

table_name = 'Ratings'
con = sqlite3.connect(filename)
cur = con.cursor()

with open(os.path.join(data_dir_ratings,filename_ratings),"rt",encoding="latin-1",newline=None) as fin: 
    contents = fin.readlines()
    contents = contents[0].split('#@#@#')
    contents = [line.replace('\'~\'', '|') for line in contents]
    columns1 = ratingsTableColumns 
    ncol     = len(columns1)
    column_list = "(" #build column names and question string
    questionstring ="("
    for n in range(ncol):
        if n+1 < (ncol):
            column_list = column_list+str(columns1[n])+", "
            questionstring = questionstring +"?," 
        else:
            column_list = column_list+str(columns1[n])+")"
            questionstring = questionstring + "?)"
    print(table_name," ",ncol)
    print(column_list)
 
    executestring = "CREATE TABLE IF NOT EXISTS "+str(table_name)+" "+column_list+";"
    cur.executescript(executestring) #create table if it does not already exist
    for row in contents:  #write remaining rows into the table
        row1=row.split('|')
        if row1[0]=='':
            row1=['' for n in range(len(columns1))]
        to_db = [str(row1[n]) for n in range(len(columns1))]         
        writestring = "INSERT INTO "+str(table_name)+" "+column_list+" VALUES "+questionstring+";"
        cur.executemany(writestring, (to_db,))
        con.commit()
con.close()

#%%
# Create auxiliary tables that link rating data to financials data
#
data_dir_crossCompanyRef = r'C:\Users\ru87118\Desktop\Compustat\Data'
filename_crossCompanyRef = r'CompanyCrossRef.txt'

crossCompanyRefTableColumns=['IdentifierId','companyId','identifierValue','identifierTypeId','startDate','endDate',
                     'activeFlag','PrimaryFlag']

table_name = 'crossCompanyRef'
con = sqlite3.connect(filename)
cur = con.cursor()

with open(os.path.join(data_dir_crossCompanyRef,filename_crossCompanyRef),encoding="utf8") as fin: 
    dr = csv.reader(fin,delimiter='|',quoting=csv.QUOTE_NONE,lineterminator='\n') #use csv.reader to read file
    columns1 = crossCompanyRefTableColumns 
    ncol     = len(columns1)
    column_list = "(" #build column names and question string
    questionstring ="("
    for n in range(ncol):
        if n+1 < (ncol):
            column_list = column_list+str(columns1[n])+", "
            questionstring = questionstring +"?," 
        else:
            column_list = column_list+str(columns1[n])+")"
            questionstring = questionstring + "?)"
    print(column_list)
    executestring = "CREATE TABLE IF NOT EXISTS "+str(table_name)+" "+column_list+";"
    cur.executescript(executestring)
    for row in dr:  
        to_db = [str(row[n]) for n in range(ncol)]         
        writestring = "INSERT INTO "+str(table_name)+" "+column_list+" VALUES "+questionstring+";"
        cur.executemany(writestring, (to_db,))
        con.commit()
con.close()

#%%
#
#
data_dir_ratingIdentifier = r'C:\Users\ru87118\Desktop\Compustat\Data'
filename_ratingIdentifier = r'spRatingIdentifier.txt'

ratingIdentifierTableColumns=['symbolId','symbolTypeId','symbolValue','objectId','relatedCompanyId','activeFlag','primaryFlag']

table_name = 'ratingIdentifier'
con = sqlite3.connect(filename)
cur = con.cursor()

with open(os.path.join(data_dir_ratingIdentifier,filename_ratingIdentifier),"rt",encoding="latin-1",newline=None) as fin: 
    contents = fin.readlines()
    contents = contents[0].split('#@#@#')
    contents = [line.replace('\'~\'', '|') for line in contents]
    columns1 = ratingIdentifierTableColumns 
    ncol     = len(columns1)
    column_list = "(" #build column names and question string
    questionstring ="("
    for n in range(ncol):
        if n+1 < (ncol):
            column_list = column_list+str(columns1[n])+", "
            questionstring = questionstring +"?," 
        else:
            column_list = column_list+str(columns1[n])+")"
            questionstring = questionstring + "?)"
    print(table_name," ",ncol)
    print(column_list)
 
    executestring = "CREATE TABLE IF NOT EXISTS "+str(table_name)+" "+column_list+";"
    cur.executescript(executestring) #create table if it does not already exist
    for row in contents:  #write remaining rows into the table
        row1=row.split('|')
        if row1[0]=='':
            row1=['' for n in range(len(columns1))]
        to_db = [str(row1[n]) for n in range(len(columns1))]         
        writestring = "INSERT INTO "+str(table_name)+" "+column_list+" VALUES "+questionstring+";"
        cur.executemany(writestring, (to_db,))
        con.commit()
con.close()

#%%
#
#
data_dir_ratingIdentifierType = r'C:\Users\ru87118\Desktop\Compustat\Data'
filename_ratingIdentifierType = r'spRatingIdentifierType.txt'

ratingIdentifierTypeTableColumns=['symbolTypeId', 'SymbolTypeName', 'objectTypeId']

table_name = 'ratingIdentifierType'
con = sqlite3.connect(filename)
cur = con.cursor()


with open(os.path.join(data_dir_ratingIdentifierType,filename_ratingIdentifierType),"rt",encoding="latin-1",newline=None) as fin: 
    contents = fin.readlines()
    contents = contents[0].split('#@#@#')
    contents = [line.replace('\'~\'', '|') for line in contents]
    columns1 = ratingIdentifierTypeTableColumns 
    ncol     = len(columns1)
    column_list = "(" #build column names and question string
    questionstring ="("
    for n in range(ncol):
        if n+1 < (ncol):
            column_list = column_list+str(columns1[n])+", "
            questionstring = questionstring +"?," 
        else:
            column_list = column_list+str(columns1[n])+")"
            questionstring = questionstring + "?)"
    print(table_name," ",ncol)
    print(column_list)
 
    executestring = "CREATE TABLE IF NOT EXISTS "+str(table_name)+" "+column_list+";"
    cur.executescript(executestring) #create table if it does not already exist
    for row in contents:  #write remaining rows into the table
        row1=row.split('|')
        if row1[0]=='':
            row1=['' for n in range(len(columns1))]
        to_db = [str(row1[n]) for n in range(len(columns1))]         
        writestring = "INSERT INTO "+str(table_name)+" "+column_list+" VALUES "+questionstring+";"
        cur.executemany(writestring, (to_db,))
        con.commit()
con.close()

#%%
# onstruct the insurance ratings table. This is the second dataset that will be used as input to python scripts
#
con = sqlite3.connect(filename)
c = con.cursor()

#c.execute("DROP TABLE CnIRatings")

with open(r'C:\Users\ru87118\Desktop\Compustat\Data\SQL\CnICompanies_create table.sql', 'r') as file:
    sql_create_CnICompanies_table = file.read().replace('\n', ' ')
    
c.execute(sql_create_CnICompanies_table)


with open(r'C:\Users\ru87118\Desktop\Compustat\Data\SQL\mapRatingsGVKEY_create table.sql', 'r') as file:
    sql_create_mapRatingsGVKEY_table = file.read().replace('\n', ' ')
    
c.execute(sql_create_mapRatingsGVKEY_table)

with open(r'C:\Users\ru87118\Desktop\Compustat\Data\SQL\CnIRatings_create table.sql', 'r') as file:
    sql_create_CnIRatings_table = file.read().replace('\n', ' ')
    
c.execute(sql_create_CnIRatings_table)

con.commit()
c.close()
con.close()

#%%
# The lines below create the dataframes used as input data
#
import pandas as pd
import os, sqlite3, csv

dir_name=r'C:\Users\ru87118\Desktop\Compustat\Data'
os.chdir(dir_name)

base_filename='CnI091920_2'
filename_suffix = '.db'
filename = os.path.join(dir_name, base_filename + filename_suffix)
con = sqlite3.connect(filename)

CnIFinancials = pd.read_sql_query("SELECT * FROM CnIFinancials", con)
CnIRatings = pd.read_sql_query("SELECT * FROM CnIRatings", con)

con.close()

CnIFinancials.to_csv(r'CnIFinancials.csv')
CnIRatings.to_csv(r'CnIRatings.csv')


# -*- coding: utf-8 -*-
"""
Created on Thu Oct  8 15:23:01 2020

@author: ru87118
"""

import pandas as pd
import os, sqlite3, csv

dir_name=r'C:\Users\ru87118\Desktop\Compustat\Data'
os.chdir(dir_name)

base_filename='CnI091920_2'
filename_suffix = '.db'
filename = os.path.join(dir_name, base_filename + filename_suffix)
con = sqlite3.connect(filename)

CnIFinancials = pd.read_sql_query("SELECT * FROM CnIFinancials", con)
CnIRatings = pd.read_sql_query("SELECT * FROM CnIRatings", con)

con.close()

#%%
CnIRatings = CnIRatings.sort_values(by=['Gvkey','ratingDate'])
CnIRatings = CnIRatings[CnIRatings['ratingTypeCode']=='STDLONG']
CnIRatings = CnIRatings.drop_duplicates(subset=['Gvkey','ratingDate'],keep='last')
CnIRatings['ratingYear'] = CnIRatings['ratingDate'].str[:4]

#Add records for continuous rating without changes
CnIRatings_ext = CnIRatings
CnIRatings_ext['ratingYear'] = CnIRatings_ext['ratingYear'].astype(int)
init = CnIRatings_ext.iloc[0]
init_year = CnIRatings_ext.iloc[0]['ratingYear']
init_key = CnIRatings_ext.iloc[0]['Gvkey']
for i in range(len(CnIRatings_ext)):
    print(i)
    if CnIRatings_ext.iloc[i]['ratingYear'] > init_year+1 and CnIRatings_ext.iloc[i]['Gvkey'] == init_key:
        for j in range(1, CnIRatings_ext.iloc[i]['ratingYear'] - init_year):
            init['ratingYear'] = init_year + j
            CnIRatings_ext = CnIRatings_ext.append(init)   
    init = CnIRatings_ext.iloc[i]
    init_year = CnIRatings_ext.iloc[i]['ratingYear']
    init_key = CnIRatings_ext.iloc[i]['Gvkey']
    
CnIRatings_ext = CnIRatings_ext.sort_values(by=['Gvkey','ratingDate'])
#%%
CnIFinancials = CnIFinancials[CnIFinancials['datafmt']=='STD']
CnIFinancials['financialYear'] = CnIFinancials['datadate'].str[:4]
CnIFinancials['financialDate'] = CnIFinancials['datadate'].str[4:]
    
CnIFinancials['financialYear'] = CnIFinancials['financialYear'].astype(int)
CnIFinancials['financialDate'] = CnIFinancials['financialDate'].astype(int)
CnIFinancials_annual = CnIFinancials[CnIFinancials['financialDate']==1231]

#%%    
CnIFinancials_annual['financialYearP1'] = CnIFinancials['financialYear'] + 1
data_merge = CnIRatings_ext.merge(CnIFinancials_annual, how='left', left_on=['Gvkey', 'ratingYear'], right_on=['GVKEY', 'financialYearP1'])

#CnIFinancials_annual.to_csv(r'CnIFinancials_annual.csv')
#CnIRatings_ext.to_csv(r'CnIRatings_ext.csv')

#data_merge = pd.read_csv(r'data_merge.csv')
#CnIFinancials_annual = pd.read_csv(r'CnIFinancials_annual.csv')
#CnIRatings_ext = pd.read_csv(r'CnIRatings_ext.csv')

data_merge_dropna = data_merge.dropna(subset=['financialYearP1'])
#data_exam = data_merge_dropna[['GVKEY', 'financialYearP1', 'Gvkey', 'ratingYear', 'ratingSymbol']]
data_merge_dropna.to_csv(r'data_merge.csv')

CnIFinancials_annual['GVKEY'].nunique()
CnIRatings['Gvkey'].nunique()
CnIRatings_ext['Gvkey'].nunique()
data_merge_dropna['GVKEY'].nunique()
print(CnIFinancials_annual.groupby('financialYear')['financialYear'].count())
print(CnIRatings_ext.groupby('ratingYear')['ratingYear'].count())
print(CnIRatings.groupby('ratingYear')['ratingYear'].count())
print(data_merge_dropna.groupby('ratingYear')['ratingYear'].count())

tmp=CnIRatings_ext.groupby('ratingYear')['ratingYear'].count()
tmp=CnIRatings.groupby('ratingYear')['ratingYear'].count()

# -*- coding: utf-8 -*-
"""
Created on Fri Nov 20 14:40:22 2015

@author: ub71894
"""


from pandas import Series, DataFrame 
import pandas as pd
import numpy as np

# -*- coding: utf-8 -*-
"""
Created on Thu Nov 19 14:34:21 2015

@author: ub71894
"""

from pandas import Series, DataFrame 
import pandas as pd
import numpy as np
from numpy.random import randn
import matplotlib.pyplot as plt

fig = plt.figure()

ax1 = fig.add_subplot(2,2,1)
ax2 = fig.add_subplot(2,2,2)
ax3 = fig.add_subplot(2,2,3)

from numpy.random import randn
plt.plot(randn(50).cumsum(),'k--')


_ = ax1.hist(randn(100), bins=20, color='k', alpha=0.3)

ax2.scatter(np.arange(30), np.arange(30) + 3 * randn(30))
#%%
fig, axes = plt.subplots(2, 2, sharex=True, sharey=True)
for i in range(2):
    for j in range(2):
        axes[i, j].hist(np.random.randn(500), bins=50, color='k', alpha=0.5)
plt.subplots_adjust(wspace=0.2, hspace=0.1)

#%% Best way
f,ax=plt.subplots()

data = randn(30).cumsum()

ax.plot(data, 'k--', label='Default')
ax.plot(data, 'k-', drawstyle='steps-post', label='steps-post')
ax.legend(loc='best')
ax.set_xlabel('dsfa')
ax.arrow(3,6)
#%%
f,ax=plt.subplots()

df = DataFrame(np.random.randn(10, 4).cumsum(0),columns=['A', 'B', 'C', 'D'],
               index=np.arange(0, 100, 10))
df.plot(ax=ax,kind='bar')

fig, axes = plt.subplots(2, 1)
data = Series(np.random.rand(16), index=list('abcdefghijklmnop'))
data.plot(kind='bar', ax=axes[0], color='k', alpha=0.7)
data.plot(kind='barh', ax=axes[1], color='k', alpha=0.7)

#%% 
import matplotlib
matplotlib.style.use('ggplot')

ts = pd.Series(np.random.randn(1000), index=pd.date_range('1/1/2000', periods=1000))
ts = ts.cumsum()

df = pd.DataFrame(np.random.randn(1000, 4), index=ts.index, columns=list('ABCD'))
df = df.cumsum()
plt.figure(); df.plot();

f,ax=plt.subplots()
comp1 = np.random.normal(0, 1, size=200) # N(0, 1)
comp2 = np.random.normal(10, 2, size=200) # N(10, 4)
values = Series(np.concatenate([comp1, comp2]))
values.hist(bins=100, alpha=0.3, color='k', normed=True,ax=ax)
values.plot(kind='kde', style='b--',ax=ax)


#%%
macro = pd.read_csv('macrodata.csv')
data = macro[['cpi', 'm1', 'tbilrate', 'unemp']]
trans_data = np.log(data).diff().dropna()


f1,ax1=plt.subplots(1,2)
f2,ax2=plt.subplots(2,1)

f1 = plt.figure(1)
plt.axes(ax1[0])
plt.scatter(trans_data['m1'], trans_data['unemp'])
plt.title('Changes in log %s vs. log %s' % ('m1', 'unemp'))
plt.axes(ax1[1])
plt.scatter(trans_data['m1'], trans_data['unemp'])
plt.title('2 second')

f2 = plt.figure(2)
plt.axes(ax2[0])
plt.scatter(trans_data['m1'], trans_data['unemp'])
plt.title('Changes in log %s vs. log %s' % ('m1', 'unemp'))
plt.axes(ax2[1])
plt.scatter(trans_data['m1'], trans_data['unemp'])
plt.title('2 second')


pd.scatter_matrix(trans_data, diagonal='kde', color='k', alpha=0.3)









# -*- coding: utf-8 -*-
"""
Created on Mon Jun 10 09:34:17 2019

@author: ub71894
"""

#%% use coin to generate Uniform (1,20) discrete number
import numpy as np
import pandas as pd
from scipy.stats import bernoulli
import seaborn as sns
import matplotlib.pyplot as plt
from math import pow

N=100000
pl_num=[]
for i in range(N):
    s=0
    for j in range(5):
        #s+= pow(2,j)*bernoulli.rvs(0.5,size=1)[0]
        s+= pow(2,j)*np.random.binomial(1,0.5) # faster
    if  s>19:
        continue
    else:
        pl_num.append(1+s)
        
    
np.array(pl_num).mean()
df= pd.DataFrame()
df['num']  = pl_num

a=df['num'].value_counts().to_frame()
a.reset_index(drop=False, inplace=True)

sns.barplot(x='index', y='num', data=a)
plt.bar(a['index'], a.num)



#%% combination function
from scipy.special import comb
comb(n,r)

import operator as op
from functools import reduce
def ncr(n, r):
    r = min(r, n-r)
    numer = reduce(op.mul, range(n, n-r, -1), 1)
    denom = reduce(op.mul, range(1, r+1), 1)
    return numer / denom



from math import factorial as fact
def ncr2(n, r):
    r = min(r, n-r)
    numer = fact(n)
    denom = fact(r)*fact(n-r)
    return numer / denom



def ncr_mul(n, pl_r):
    numer = fact(n)
    denom = 1
    for num in pl_r:
        denom *= fact(num)        
    return numer / denom



#%% 
import numpy as np
from numba import jit
@jit
def sim(N):
    pl_result =[]
    for i in range(N):
        a = np.random.randint(1,10001,100)
        b = np.random.randint(1,10001,100)
        s=0
        for num in b:
            if num in a:
                s+=1
        pl_result.append(s)
    return (pl_result)


np.array(sim(500000)).mean()
sim(100)

#%% shuffle 1,2,...N and in the shuffled list, what's the expect number of pair
#(n,n+1) that list[n+1] = list[n]+1, answer = (N-1)/N
 
@jit
def simu(N):
    arr = np.arange(N)
    np.random.shuffle(arr)
    Y=0
    for i in range (N-1):
        if (arr[i+1] == arr[i]+1):
            Y+=1
    return(Y)
s=0
N=10000000
for i in range(N):
   s+=simu(5)
s/N
#%% plot random number 
import matplotlib.pyplot as plt
y = np.random.rand(1000)
x = np.arange(1000)

fig, ax = plt.subplots()
ax.scatter(x,y)


#%% 10 normal rv's max's exp
@jit
def simu(N):
    s=0
    for i in range(N):
        x = np.random.randn(10)
        s+=max(x)
    return (s/N)


#%% 
from scipy.stats import norm
norm.cdf(2)
norm.ppf(.975)
# 1.959963984540054
norm.ppf(.975, 1, 2)
# 4.919927969080108

from scipy.stats import f
f.ppf(0.975, 32, 59)
# 1.8022087474311501
f.cdf(2, 32, 59)

from scipy.stats import chi2
df=5
x = np.linspace(chi2.ppf(0.01, df),chi2.ppf(0.99, df), 100)
fig, ax = plt.subplots()
ax.plot(x, chi2.pdf(x, df),'r-', lw=5, alpha=0.6, label='chi2 pdf')


from scipy.stats import t
df = 30
mean, var, skew, kurt = t.stats(df, moments='mvsk')

t.ppf(0.975, 10)
#2.2281388519649385
#%% 
dat
'''
Out[147]: 
     id  month  times
0    yu      1    100
1  wang      1    100
2  wang      2    200
3  wang      3    300
4   sun      1    100
5   sun      2    200
6   sun      3    300
7   sun      4    400
'''
aa= dat.groupby(by=['id']).mean()['times']

#%% Fib

# Recurssive program to find n'th fibonacci number 
def fib(n): 
    if n <= 1: 
        return n 
    return fib(n-1) + fib(n-2) 

#%% 2 sample test
from scipy.stats import ttest_ind
a = np.random.randn(1000)+0.5
b = np.random.randn(1000)+0.4
ttest_ind(a,b)  
ttest_ind(a,b, equal_var=False)    

# rank sum / Wilcoxon test
from scipy.stats import ranksums
ranksums(a,b)

from scipy.stats import wilcoxon  # signed rank test
a = np.random.randn(1000)+0.5
b = np.random.randn(1000)+0.4
c=a-b
wilcoxon(c)

from scipy.stats import ks_2samp
ks_2samp(a, b)# power is low sometime



#%% other tests
from scipy.stats import levene
a = np.random.randn(1000)+0.3
b = np.random.randn(1000)+0.4
levene(a,b)
from scipy.stats import f
F = a.var() / b.var()
f.cdf(F, 999,999)
from scipy.stats import f_oneway
f_oneway(a,b)  # equilevent with t test

from scipy.stats import chisquare
chisquare([16, 18, 16, 14, 12, 12]) # assume uniform dist
chisquare([16, 18, 16, 14, 12, 12], f_exp=[14.66,]*6)
chisquare([16, 18, 16, 14, 12, 12], f_exp=[16, 16, 16, 16, 16, 8])


# normality test
from scipy.stats import kstest, shapiro, normaltest
a = np.random.randn(1000)
kstest(a, 'norm')
kstest(a, 't', args=[10])
shapiro(a)
normaltest(a) # D’Agostino




#%% 
import statsmodels.api as sm
import numpy as np
import pandas as pd
df = pd.DataFrame()
df['x1'] = np.random.rand(1000)
df['x2']= 3*np.random.randn(1000)

y = 5*df.x1 + 2*df.x2 + np.random.randn(1000)
df['y']=y
X = sm.add_constant(df[['x1','x2']])
ss=sm.OLS(y, X).fit()
ss.summary()


#%% ANOVA
#http://www.statsmodels.org/devel/generated/statsmodels.stats.anova.anova_lm.html


#%% truncated normal dist
import numpy as np
from scipy.stats import norm
from scipy.optimize import root

mu = 1
sigma = 4
x= np.random.normal(mu, sigma,100000)

T = -5
y = x[x<T]

Ex = y.mean()
EVar = y.var()

phi = norm.pdf
PHI = norm.cdf


def fun1(mu, sigma):
    beta = (T- mu)/sigma
    return (mu-sigma*phi(beta)/PHI(beta) - Ex )
def fun2(mu, sigma):
    beta = (T- mu)/sigma
    return (sigma**2*(1-beta*phi(beta)/PHI(beta)-phi(beta)**2/PHI(beta)**2) - EVar)

def ssfun(x):
    return ([fun1(x[0],x[1]), fun2(x[0],x[1])])

root(ssfun, [0.5,0.5])




mu=1
sigma=4
beta = (T- mu)/sigma
mu-sigma*phi(beta)/PHI(beta)
sigma**2*(1-beta*phi(beta)/PHI(beta)-phi(beta)**2/PHI(beta)**2)

Ex
EVar
# -*- coding: utf-8 -*-
"""
Created on Fri Aug 14 15:53:13 2015

@author: ub69401
"""

import os,sqlite3

os.chdir(r"C:\Users\ub69401\Documents\compustat\database") #change this to the directory where your csv files are stored
txt_files={} # we store the dataframes in a dictionary

dir_name=r'C:\Users\ub69401\Documents\compustat\database'
base_filename='compustat'
filename_suffix = '.db'
filename = os.path.join(dir_name, base_filename + filename_suffix)
sqlite_file = filename
con = sqlite3.connect(filename)
cur = con.cursor()

COMPUSTAT = "'USA'"
executestring = "SELECT * FROM COMPANY WHERE LOC="+str(COMPUSTAT)

cur.execute(executestring)

rows = cur.fetchall()
for row in rows:
    print(row)
    
con.close()# -*- coding: utf-8 -*-
"""
Created on Fri Aug 14 15:53:13 2015

@author: ub71894
"""

import os,sqlite3
import pandas as pd

os.chdir(r"C:\Users\ub71894\Documents\code\Python") #change this to the directory where your csv files are stored
txt_files={} # we store the dataframes in a dictionary

dir_name=r'C:\Users\ub71894\Documents\data\REIG_TOPS\db'
base_filename='CMO South'
filename_suffix = '.db'
filename = os.path.join(dir_name, base_filename + filename_suffix)
con = sqlite3.connect(filename)
cur = con.cursor()

executestring = 'SELECT * FROM CMD_PROPERTY_INFO_HISTORY'

cur.execute(executestring)

names = list(map(lambda x: x[0], cur.description))
data=pd.DataFrame(cur.fetchall(),columns=names)




#rows = cur.fetchall()
#for row in rows:
#    print(row)
    
con.close()



import os, sqlite3, csv
from io import StringIO


dir_name=r'C:\Users\ru87118\Desktop\Compustat\Data'
os.chdir(dir_name)

data_dir = r'C:\Users\ru87118\Desktop\Compustat\Data'
filename_A_L = r'f_co_afnd1V2.txt'
remove_string_A_L = 'H|20200919-07:42:53|co_afnd1V2|6|'

file_A_L = open(os.path.join(data_dir, filename_A_L),'r').read().replace(remove_string_A_L,'').replace('|',',')
s = StringIO(file_A_L)
# Save a csv for direct manual inspection
with open('afnd1V2.csv', 'w') as f:
    for line in s:
        #print(line)
        f.write(line)
        

filename_M_Z = r'f_co_afnd2V2.txt'
remove_string_M_Z = 'H|20200919-07:41:52|co_afnd2V2|6|'

file_M_Z = open(os.path.join(data_dir, filename_M_Z),'r').read().replace(remove_string_M_Z,'').replace('|',',')
s = StringIO(file_M_Z)
# Save a csv for direct manual inspection
with open('afnd2V2.csv', 'w') as f:
    for line in s:
        f.write(line)
        

filename_company = r'f_company.txt'
remove_string_company = 'H|20200919-07:32:36|company|1|'

file_company = open(os.path.join(data_dir, filename_company),'r').read().replace(remove_string_company,'').replace('|',',')
s = StringIO(file_company)
# Save a csv for direct manual inspection
with open('company.csv', 'w') as f:
    for line in s:
        f.write(line)        
        
        
filename_rating = r'spRatingData.txt'

file_rating = open(os.path.join(data_dir, filename_rating),'r').read().replace('\'~\'',',').replace('#@#@#','\n')
s = StringIO(file_rating)

ratingsTableColumns = str(['ratingDetailId','entitySymbolValue','instrumentSymbolValue','securitySymbolValue',
                     'objectTypeId'	,'orgDebtTypeCode','ratingTypeCode','currentRatingSymbol','ratingSymbol',
                     'ratingDate','creditwatch','outlook','creditwatchDate','outlookDate','priorRatingSymbol',	
                     'priorCreditwatch','priorOutlook','ratingQualifier','regulatoryIndicator','regulatoryQualifier',
                     'ratingActionWord','CWOLActionWord','cwolInd','maturityDate','CUSIP','CINS','ISIN'])

# Save a csv for direct manual inspection
with open('spRatingData.csv', 'w') as f:
    f.write(ratingsTableColumns) 
    f.write("\n")
    for line in s:
        f.write(line)        
        

filename_key = r'CompanyCrossRef.txt'

file_key = open(os.path.join(data_dir, filename_key),'r').read().replace('|',',')
s = StringIO(file_key)
# Save a csv for direct manual inspection
with open('CompanyCrossRef.csv', 'w') as f:
    for line in s:
        f.write(line)             
        
        
import pandas as pd
import numpy as np 
#import scipy as sp 
#from sklearn.linear_model import LogisticRegression, LinearRegression
import statsmodels.api as sm 
import seaborn as sns
#import statsmodels.formula.api as smf


x = np.random.randn(1000,3)
alpha=4.5; beta = pd.Series([11.3,4.23])
x=pd.DataFrame(x)
y_true = alpha + x[0]*beta[0] + x[1]*beta[1] 
y = y_true + x[2]

inde = sm.add_constant(x[[0,1]],prepend = True)

mod = sm.OLS(y, inde)
res = mod.fit()

ybar = y.mean()
SST = sum((y-ybar)**2)
SSE = sum((y-res.fittedvalues)**2)

SSR = sum((ybar-res.fittedvalues)**2)

res.summary()

dat = pd.concat([y,inde],axis=1,names='haha')
dat.columns = ['ysome','const','x1','x2']

sns.regplot(x="x1",y="ysome",data=dat)
#%%
import matplotlib.pyplot as plt
from statsmodels.sandbox.regression.predstd import wls_prediction_std


prstd, iv_l, iv_u = wls_prediction_std(res)

fig, ax = plt.subplots(figsize=(8,6))

ax.plot(x[0], y, 'o', label="data")
ax.plot(x[0], y_true, 'bo', label="True")
ax.plot(x[0], res.fittedvalues, 'ro', label="OLS")
ax.plot(x[0], iv_u, 'ro')
ax.plot(x[0], iv_l, 'ro')
ax.legend(loc='best');




#%% experiment:

x = np.random.randn(1000,3)
x[0] =  x[0]*3+2; x[1] =  x[1]*2+5;

alpha=14.5; beta = pd.Series([11.3,4.23])
x=pd.DataFrame(x)
y_true = alpha + x[0]*beta[0] + x[1]*beta[1] 
y = y_true + x[2]

# with intercept:
inde = sm.add_constant(x[[0,1]],prepend = True)
mod = sm.OLS(y, inde)
res = mod.fit()
res.summary()



#without intercept
mod = sm.OLS(y, x[[0,1]])
res = mod.fit()
res.summary()

xx=x.copy()
xx[0]= x[0]*res.params[0]
xx[1]= x[1]*res.params[1]

inde = sm.add_constant(xx[[0,1]],prepend = True)
mod = sm.OLS(y, inde)
res = mod.fit()
res.summary()


x12= x[0]*res.params[0]+x[1]*res.params[1]
inde = sm.add_constant(x12,prepend = True)
mod = sm.OLS(y, inde)
res = mod.fit()

res.summary()




#%%
plt.subplot(221)
plt.scatter(dat.Totalscore, dat.Final_PD_Risk_Rating, s=25)    
plt.subplot(224)
plt.scatter(dat.Totalscore, dat.Prelim_PD_Risk_Rating_Uncap, s=25)    


# works
plt.subplot(221)
sns.scatterplot(x='quantscore', y='Prelim_PD_Risk_Rating_Uncap', data=dat)
plt.subplot(224)
sns.scatterplot(x='quantscore', y='Final_PD_Risk_Rating', data=dat)


# works
ax=plt.subplot(221)
sns.regplot(x='quantscore', y='Prelim_PD_Risk_Rating_Uncap', data=dat)
plt.subplot(224)
plt.scatter(dat.quantscore, dat.Prelim_PD_Risk_Rating_Uncap, s=25)    


# works
fig, ax = plt.subplots(3,3, sharex='col')
sns.regplot(x='quantscore', y='Prelim_PD_Risk_Rating_Uncap', data=dat,ax=ax[0,1])
ax[2,1].scatter(dat.quantscore, dat.Prelim_PD_Risk_Rating_Uncap, s=25) 
from IPython import get_ipython
get_ipython().magic('reset -f')
get_ipython().magic('clear')import os,sqlite3
from sqlite3 import OperationalError
from sqlite3 import ProgrammingError

os.chdir(r"C:\Users\ub71894\Documents\Python Scripts\sql")
sql_file = 'liteworld.sql'
con = sqlite3.connect('world2.db')
cursor = con.cursor()


def exec_sql_file(cursor, sql_file):
    print ("\n[INFO] Executing SQL script file: '%s'" % (sql_file))
    statement = ""

    for line in open(sql_file):
        if re.match(r'--', line):  # ignore sql comment lines
            continue
        if not re.search(r'[^-;]+;', line):  # keep appending lines that don't end in ';'
            statement = statement + line
        else:  # when you get a line ending in ';' then exec statement and reset for next statement
            statement = statement + line
            #print "\n\n[DEBUG] Executing SQL statement:\n%s" % (statement)
            try:
                cursor.execute(statement)
            except (OperationalError, ProgrammingError) as e:
                print ("\n[WARN] MySQLError during execute statement \n\tArgs: '%s'" % (str(e.args)))

            statement = ""    


exec_sql_file(cursor, sql_file)

con.commit()

con.close()

# -*- coding: utf-8 -*-
"""
Spyder Editor
http://docs.scipy.org/doc/scipy/reference/tutorial/optimize.html
This is a temporary script file.
"""

import numpy as np
from scipy.optimize import minimize


def func(x, sign=1.0):
    """ Objective function """
    return sign*(2*x[0]*x[1] + 2*x[0] - x[0]**2 - 2*x[1]**2)


def func_deriv(x, sign=1.0):
    """ Derivative of objective function """
    dfdx0 = sign*(-2*x[0] + 2*x[1] + 2)
    dfdx1 = sign*(2*x[0] - 4*x[1])
    return np.array([ dfdx0, dfdx1 ])

cons = ({'type': 'eq','fun' : lambda x: np.array([x[0]**3 - x[1]]),
    'jac' : lambda x: np.array([3.0*(x[0]**2.0), -1.0])},
    {'type': 'ineq', 'fun' : lambda x: np.array([x[1] - 1]),
    'jac' : lambda x: np.array([0.0, 1.0])})


res = minimize(func, [-1.0,1.0], args=(-1.0,), jac=func_deriv, constraints=cons, 
               method='SLSQP', options={'disp': True})
               
print(res.x)               

#%%
from scipy.interpolate import interp1d
x = np.linspace(0, 10, num=11, endpoint=True)
y = np.cos(-x**2/9.0)
f = interp1d(x, y)
f2 = interp1d(x, y, kind='cubic')

xnew = np.linspace(0, 10, num=41, endpoint=True)
import matplotlib.pyplot as plt
plt.plot(x, y, 'o', xnew, f(xnew), '-', xnew, f2(xnew), '--')
plt.legend(['data', 'linear', 'cubic'], loc='best')
plt.show()


#%%
def func(x, y):
    return x*(1-x)*np.cos(4*np.pi*x) * np.sin(4*np.pi*y**2)**2


grid_x, grid_y = np.mgrid[0:1:100j, 0:1:200j]
points = np.random.rand(1000, 2)
values = func(points[:,0], points[:,1])
#%%
from scipy.interpolate import griddata
grid_z0 = griddata(points, values, (grid_x, grid_y), method='nearest')
grid_z1 = griddata(points, values, (grid_x, grid_y), method='linear')
grid_z2 = griddata(points, values, (grid_x, grid_y), method='cubic')

import matplotlib.pyplot as plt
plt.subplot(221)
plt.imshow(func(grid_x, grid_y).T, extent=(0,1,0,1), origin='lower')
plt.plot(points[:,0], points[:,1], 'k.', ms=1)
plt.title('Original')
plt.subplot(222)
plt.imshow(grid_z0.T, extent=(0,1,0,1), origin='lower')
plt.title('Nearest')
plt.subplot(223)
plt.imshow(grid_z1.T, extent=(0,1,0,1), origin='lower')
plt.title('Linear')
plt.subplot(224)
plt.imshow(grid_z2.T, extent=(0,1,0,1), origin='lower')
plt.title('Cubic')
plt.gcf().set_size_inches(6, 6)
plt.show()

#%%

from scipy import stats
from scipy.stats import norm

norm.rvs(loc=10, scale=5, size=(4,3))# -*- coding: utf-8 -*-
"""
Created on Thu Nov 12 11:26:50 2015

@author: ub71894
"""

import numpy as np
import pandas as pd
from scipy import stats, integrate
import matplotlib.pyplot as plt
import seaborn as sns
import string




sns.set(style="white")

# Generate a large random dataset
rs = np.random.RandomState(33)
d = pd.DataFrame(data=rs.normal(size=(100, 26)),
                 columns=list(string.ascii_letters[:26]))

# Compute the correlation matrix
corr = d.corr()

# Generate a mask for the upper triangle
mask = np.zeros_like(corr, dtype=np.bool)
mask[np.triu_indices_from(mask)] = True

# Set up the matplotlib figure
f, ax = plt.subplots(figsize=(11, 9))

# Generate a custom diverging colormap
cmap = sns.diverging_palette(220, 10, as_cmap=True)

# Draw the heatmap with the mask and correct aspect ratio
sns.heatmap(corr, mask=mask, cmap=cmap, vmax=.3,
            square=True, xticklabels=5, yticklabels=5,
            linewidths=.5, cbar_kws={"shrink": .5}, ax=ax)





import seaborn as sns
sns.set()
#fig1
seafig= sns.pairplot(trans_data)

#fig2
seafig = sns.pairplot(trans_data, palette="Set2", diag_kind="kde", size=2.5)

#fig3
seafig = sns.PairGrid(trans_data)
seafig.map_upper(plt.scatter)
seafig.map_lower(sns.kdeplot, cmap="Blues_d")
seafig.map_diag(sns.kdeplot, lw=3, legend=False);

seafig.savefig('seaborn.pdf')





sns.set()
plt.figure(figsize=(18, 6))
# Create a random dataset across several variables
rs = np.random.RandomState(0)
n, p = 40, 12
d = rs.normal(0, 2, (n, p))
d += np.log(np.arange(1, p + 1)) * -5 + 10

# Use cubehelix to get a custom sequential palette
pal = sns.cubehelix_palette(p, rot=-.5, dark=.3)

# Show each distribution with both violins and points
d= pd.DataFrame(d,columns=list('abcdefghijkl'))
sns.violinplot(data=d, palette=pal, inner="quart")

plt.savefig('seaborn.pdf')


import os
os.chdir(r"C:\Users\ub71894\Documents\Data")
import pandas as pd 
import numpy as np

data = pd.read_excel(r'ind dev quant.xls')

data2 = pd.read_excel(r'ind dev quant.xls')

quant_factor = ['quant1','quant2','quant3','quant4']

floor = [0.024234, 1.84363, 76699.5, 15196]
cap = [2.33655, 406.7364, 36700450, 5856121]
  


# cap/floor:
for i, col in enumerate(quant_factor):
    data[col] = np.clip(data[col], floor[i], cap[i])


#data['quant1'].fillna(data['quant1'].mean(), inplace=True)

data[quant_factor].describe()



len(data.query('quant1!=quant1'))
len(data.query('quant2!=quant2'))
len(data.query('quant3!=quant3'))
len(data.query('quant4!=quant4'))

len(data.query('quant1!=quant1 or quant2!=quant2 '))

len(data.query('quant1!=quant1 or quant3!=quant3 '))

len(data.query('quant1!=quant1 or quant4!=quant4 '))

len(data.query('quant2!=quant2 or quant3!=quant3 '))

len(data.query('quant4!=quant4 or quant3!=quant3 '))















#%%

data = pd.read_excel(r'ind dev quant.xls')
data.dropna(subset=['quant4','quant3'], how='any', inplace=True)

data.loc[data.quant2<0,"quant2"]=np.nan

# cap/floor:
for i, col in enumerate(quant_factor):
    data[col] = np.clip(data[col], floor[i], cap[i])

data.quant2.fillna(406.7364, inplace=True)
data['quant1'].fillna(data['quant1'].mean(), inplace=True)



data[quant_factor].describe()


def ssum(list, sum):
    current = ""
    ssum_h(list, len(list), current, sum)

def ssum_h(list, n, subset, sum):
    if sum == 0:
        print (subset)
        return
      
    if n == 0:
        return
      
    if list[n-1] <= sum:
        ssum_h(list, n-1, subset, sum)
        ssum_h(list, n-1, subset+str(list[n-1])+" ", sum-list[n-1])
    else:
        ssum_h(list, n-1, subset, sum)


lst=[280072,14676,286875,1690,762,1148,4815,7704,15408,14445,8828,8581,54021,50021,18329,7338,890894,45645]
sum=328519
ssum(lst, sum)
'''
                                       
      3.141592653589793238462643383279  
    5028841971693993751058209749445923  
   07816406286208998628034825342117067  
   9821    48086         5132           
  823      06647        09384           
 46        09550        58223           
 17        25359        4081            
           2848         1117            
           4502         8410            
           2701         9385            
          21105        55964            
          46229        48954            
          9303         81964            
          4288         10975            
         66593         34461            
        284756         48233            
        78678          31652        71  
       2019091         456485       66  
      9234603           486104543266485
     2133936            0726024914127   
     3724587             00660631558    
     817488               152092096     
'''

from dirsync import sync

#%% from local to I drive
local = [
r'C:\Users\ub71894\Documents\Projects\CNI',
r'C:\Users\ub71894\Documents\Projects\MachineLearning',
r'C:\Users\ub71894\Documents\Data',
r'C:\Users\ub71894\Documents\Code',
r'C:\Users\ub71894\Documents\DevRepo',
r'C:\Users\ub71894\Documents\Projects\AnnualPerformance',
r'C:\Users\ub71894\Documents\Projects\GlobalCorporation',
r'C:\Users\ub71894\Documents\Projects\CNI_redev'
]

Udrive = [
r'I:\YangY\Projects\CNI', 
r'I:\YangY\Projects\MachineLearning', 
r'I:\YangY\Data', 
r'I:\YangY\Code', 
r'I:\YangY\DevRepo',
r'I:\YangY\Projects\AnnualPerformance',
r'I:\YangY\Projects\GlobalCorporation',
r'I:\YangY\Projects\CNI_redev'
]

for i in range(len(local)):
	sync(local[i], Udrive[i], action='sync', purge=True, create=True)

print("I: drive sync finished")
#%% from local to U drive
local = [
r'C:\Users\ub71894\Documents\Code', 
r'C:\Users\ub71894\Documents\Docs', 
r'C:\Users\ub71894\Documents\Projects',
#r'C:\Users\ub71894\Documents\Data',
#r'C:\Users\ub71894\Documents\DevRepo',
r'C:\Users\ub71894\Documents\OneNote Notebooks'
]

Udrive = [
r'U:\Code', 
r'U:\Docs', 
r'U:\Projects', 
#r'U:\Data', 
#r'U:\DevRepo', 
r'U:\OneNote Notebooks']

for i in range(len(local)):
	sync(local[i], Udrive[i], action='sync', purge=True, create=True)




import pandas as pd
import numpy as np
from numba import jit
from matplotlib import pyplot as plt


N = 30
N_score = np.array(score*N)
assign_score = np.random.choice(N_score, N*4, replace=False)
ind_score = np.zeros(N)
for j in range(N):
    ind_score[j] = assign_score[j*4:j*4+4].mean()

ind_score_ranked = ind_score[ind_score.argsort()]


@jit
def simulation(N, N_simu, rank):
    score = [1,2,3,4]
    N_score = np.array(score*N)
    result = np.zeros(N_simu)
    for i in range(N_simu):
        assign_score = np.random.choice(N_score, N*4, replace=False)
        ind_score = np.zeros(N)
        for j in range(N):
            ind_score[j] = assign_score[j*4:j*4+4].mean()
        ind_score_ranked = ind_score[ind_score.argsort()]
        result[i] = ind_score_ranked[-rank]

    return (result)


def plot_rank(rank = 1):


    aa1 = simulation(3, 100000, rank)
    aa2 = simulation(4, 100000, rank)
    aa3 = simulation(6, 100000, rank)
    aa4 = simulation(7, 100000, rank)
    aa5 = simulation(30, 100000, rank)
    
    
    ax = plt.subplot(511)
    ax.set_autoscaley_on(False)
    
    plt.hist(aa1, histtype='stepfilled', bins=30, alpha=0.85,
             label=f"""Team_3 std={aa1.std()},
             mean={aa1.mean()} """, color="#A60628", normed=True)
    plt.legend(loc="upper left")
    plt.title(f"some interesting analysis")
    plt.xlim([1, 4])
    plt.xlabel(f"mean score of rank {rank}")
    
    ax = plt.subplot(512)
    ax.set_autoscaley_on(False)
    plt.hist(aa2, histtype='stepfilled', bins=30, alpha=0.85,
             label=f"""Team_4 std={aa2.std()},
             mean={aa2.mean()} """, color="#7A68A6", normed=True)
    plt.legend(loc="upper left")
    plt.xlim([1, 4])
    plt.xlabel(f"mean score of rank {rank}")
    
    
    ax = plt.subplot(513)
    ax.set_autoscaley_on(False)
    plt.hist(aa3, histtype='stepfilled', bins=30, alpha=0.85,
             label=f"""Team_6 std={aa3.std()},
             mean={aa3.mean()} """, color="#1A68A6", normed=True)
    plt.legend(loc="upper left")
    plt.xlim([1, 4])
    plt.xlabel(f"mean score of rank {rank}")
    
    ax = plt.subplot(514)
    ax.set_autoscaley_on(False)
    plt.hist(aa4, histtype='stepfilled', bins=30, alpha=0.85,
             label=f"""Team_7 std={aa4.std()},
             mean={aa4.mean()} """, color="#1A68A6", normed=True)
    plt.legend(loc="upper left")
    plt.xlim([1, 4])
    plt.xlabel(f"mean score of rank {rank}")
    
    
    ax = plt.subplot(515)
    ax.set_autoscaley_on(False)
    plt.hist(aa5, histtype='stepfilled', bins=30, alpha=0.85,
             label=f"""Team_30 std={aa5.std()},
             mean={aa5.mean()} """, color="#1A68A6", normed=True)
    plt.legend(loc="upper left")
    plt.xlim([1, 4])
    plt.xlabel(f"mean score of rank {rank}")
# -*- coding: utf-8 -*-
"""
Created on Fri Feb 26 13:45:55 2016

@author: ub71894 (4e8e6d0b), CSG
"""


class MyObject(object):
    def __init__(self,x,y):
        self._x = x
        self._y = y

    @property
    def x(self):
        return self._x

    @x.setter
    def x(self, value):
        self._x = value

    @property
    def y(self):
        return self._y

    def power(self):
        return self._x * self._y

aa= MyObject(4,9)

aa.power()
aa.y
aa.x

aa.x=10
aa.y=12

#%%
import numpy as np
from pandas import Series, DataFrame
import pandas as pd
from datetime import datetime

long = Series(np.random.randn(1000).cumsum(), index=pd.date_range('1/1/2000', periods=1000))
long.plot()
pd.rolling_mean(long,50).plot()



std50 = pd.rolling_std(long,50)
std50.plot()'''original example for checking how far GAM works
Note: uncomment plt.show() to display graphs
'''

example = 3  # 1,2 or 3

import numpy as np
import numpy.random as R
import matplotlib.pyplot as plt

from statsmodels.sandbox.gam import AdditiveModel
from statsmodels.sandbox.gam import Model as GAM #?
from statsmodels.genmod.families import family
from statsmodels.genmod.generalized_linear_model import GLM

standardize = lambda x: (x - x.mean()) / x.std()
demean = lambda x: (x - x.mean())
nobs = 150
x1 = R.standard_normal(nobs)
x1.sort()
x2 = R.standard_normal(nobs)
x2.sort()
y = R.standard_normal((nobs,))

f1 = lambda x1: (x1 + x1**2 - 3 - 1 * x1**3 + 0.1 * np.exp(-x1/4.))
f2 = lambda x2: (x2 + x2**2 - 0.1 * np.exp(x2/4.))
z = standardize(f1(x1)) + standardize(f2(x2))
#z = standardize(z) * 2 # 0.1

y += z
d = np.array([x1,x2]).T


if example == 1:
    print("normal")
    m = AdditiveModel(d)
    m.fit(y)
    x = np.linspace(-2,2,50)

    print(m)

    y_pred = m.results.predict(d)
    plt.figure()
    plt.plot(y, '.')
    plt.plot(z, 'b-', label='true')
    plt.plot(y_pred, 'r-', label='AdditiveModel')
    plt.legend()
    plt.title('gam.AdditiveModel')

import scipy.stats, time

if example == 2:
    print("binomial")
    f = family.Binomial()
    b = np.asarray([scipy.stats.bernoulli.rvs(p) for p in f.link.inverse(y)])
    b.shape = y.shape
    m = GAM(b, d, family=f)
    toc = time.time()
    m.fit()
    tic = time.time()
    print(tic-toc)


if example == 3:
    print("Poisson")
    f = family.Poisson()
    y = y/y.max() * 3
    yp = f.link.inverse(y)
    p = np.asarray([scipy.stats.poisson.rvs(p) for p in f.link.inverse(y)], float)
    p.shape = y.shape
    m = GAM(p, d, family=f)
    toc = time.time()
    m.fit(p)
    tic = time.time()
    print(tic-toc)


plt.figure()
plt.plot(x1, standardize(m.smoothers[0](x1)), 'r')
plt.plot(x1, standardize(f1(x1)), linewidth=2)
plt.figure()
plt.plot(x2, standardize(m.smoothers[1](x2)), 'r')
plt.plot(x2, standardize(f2(x2)), linewidth=2)




plt.show()import os
os.environ['PATH'] += os.pathsep + r'C:\Users\ub71894\AppData\Graphviz2.38'
os.environ['PATH'] += os.pathsep + r'C:\Users\ub71894\AppData\Graphviz2.38\bin'
from graphviz import Digraph
dot = Digraph(comment='The Round Table')
dot.node('A', 'King Arthur')
dot.node('B', 'Sir Bedevere the Wise')
dot.node('L', 'Sir Lancelot the Brave')
dot.edges(['AB', 'AL'])
dot.edge('B', 'L', constraint='false')


dot
import pandas as pd
import numpy as np
from sklearn import linear_model
from sklearn import datasets
diabetes= datasets.load_diabetes()
dat= diabetes['data']
y = diabetes['target']>140

#%%
model = linear_model.LogisticRegression(fit_intercept=True, penalty='l2', dual=False, tol=0.0001, C=1000000000)
model.fit(dat,y)
model.coef_
model.intercept_

#%%
import statsmodels.api as sm 
dat2 = pd.DataFrame(dat, columns=['a','b','c','d','e','f','g','h','i','j'])
dat2['y'] = y 
x=sm.add_constant(dat2[['a','b','c','d','e','f','g','h','i','j']],prepend = True)
y=dat2['y']

logit = sm.Logit(y, x)
result = logit.fit()
print (result.summary())



#%%
X = dat.T;  Y = y.reshape((1,442))              
W_new = np.random.randn(10,1);    b_new=0
W_old = W_new+1;    b_old = 1

alpha=0.01 / 442
def sigmoid(z):
    return(1/(1+np.exp(-z)))
n=0
while np.linalg.norm(W_new-W_old)>0.0001:
    
#while n <= 10:
    W_old = W_new;      b_old = b_new; 
    A = sigmoid(np.dot(W_old.T,X)+b_old)
    W_new = W_old - alpha*np.dot(X,(A-Y).T)

    b_new = b_old - alpha*(A-Y).T.sum()

    cost = -np.dot(Y,np.log(A).T)-np.dot((1-Y),np.log(1-A).T)
    print("after "+str(n)+" iterations, the cost function is "+ str(cost))
    n +=1


W_new = np.array([1,-11,13,11,-32,22,-0.5,1,30,0]); W_new=W_new.reshape((10,1));    b_new=0
W_old = W_new+1;    b_old = 1
n=0
while np.linalg.norm(W_new-W_old)>0.0001:
    
#while n <= 10:
    W_old = W_new;      b_old = b_new; 
    A = sigmoid(np.dot(W_old.T,X)+b_old)
    W_new = W_old - alpha*np.dot(X,(A-Y).T)

    b_new = b_old - alpha*(A-Y).T.sum()

    cost = -np.dot(Y,np.log(A).T)-np.dot((1-Y),np.log(1-A).T)
    print("after "+str(n)+" iterations, the cost function is "+ str(cost))
    n +=1




W_new = np.zeros((10,1));    b_new=0
W_old = W_new+1;    b_old = 1
n=0
while np.linalg.norm(W_new-W_old)>0.0001:
    
#while n <= 10:
    W_old = W_new;      b_old = b_new; 
    A = sigmoid(np.dot(W_old.T,X)+b_old)
    W_new = W_old - alpha*np.dot(X,(A-Y).T)

    b_new = b_old - alpha*(A-Y).T.sum()

    cost = -np.dot(Y,np.log(A).T)-np.dot((1-Y),np.log(1-A).T)
    print("after "+str(n)+" iterations, the cost function is "+ str(cost))
    n +=1



#%% SGD
W_new = np.zeros((10,1));    b_new=0
W_new = np.array([1,-11,13,11,-32,22,-0.5,1,30,0]); W_new=W_new.reshape((10,1));    b_new=0
W_new = np.random.randn(10,1);    b_new=0

W_old = W_new+1;    b_old = 1
n=0; alpha=0.01

while np.linalg.norm(W_new-W_old)>0.00001:
    
    for i in range(442):       

        X_i = X[:,i].reshape((10,1)); Y_i = Y[0,i]
        W_old = W_new;      b_old = b_new; 
        A_i = sigmoid(np.dot(W_old.T,X_i)+b_old)
        W_new = W_old - alpha*X_i*(A_i-Y_i) 

        b_new = b_old - alpha*(A_i-Y_i).sum()

    A = sigmoid(np.dot(W_new.T,X)+b_new)
    cost = -np.dot(Y,np.log(A).T)-np.dot((1-Y),np.log(1-A).T)
    print("after "+str(n)+" iterations, the cost function is "+ str(cost))
    n +=1

import pandas as pd
import statsmodels.api as sm 
import statsmodels.formula.api as smf

data = pd.read_table('CHDAGE.txt', sep='\s+')


#%% Using statsmodels.api.Logit
dat=data.copy()
dat['const'] = 1
x=sm.add_constant(dat['AGE'],prepend = True)
y=(dat.CHD)

logit = sm.Logit(y, x)
result = logit.fit()
print (result.summary())

'''            Logit Regression Results                           
==============================================================================
Dep. Variable:                    CHD   No. Observations:                  100
Model:                          Logit   Df Residuals:                       98
Method:                           MLE   Df Model:                            1
Date:                Fri, 13 Nov 2015   Pseudo R-squ.:                  0.2145
Time:                        11:24:53   Log-Likelihood:                -53.677
converged:                       True   LL-Null:                       -68.331
                                        LLR p-value:                 6.168e-08
==============================================================================
                 coef    std err          z      P>|z|      [95.0% Conf. Int.]
------------------------------------------------------------------------------
const         -5.3095      1.134     -4.683      0.000        -7.531    -3.088
AGE            0.1109      0.024      4.610      0.000         0.064     0.158
==============================================================================
'''




#%% Using statsmodels.api.GLM

dat=data.copy()
dat['const'] = 1
x=sm.add_constant(dat['AGE'],prepend = True)
y=(dat.CHD)


glm_lr = sm.GLM(y, x, family=sm.families.Binomial())
glm_results = glm_lr.fit()
print(glm_results.summary())


'''
                 Generalized Linear Model Regression Results                  
==============================================================================
Dep. Variable:                    CHD   No. Observations:                  100
Model:                            GLM   Df Residuals:                       98
Model Family:                Binomial   Df Model:                            1
Link Function:                  logit   Scale:                             1.0
Method:                          IRLS   Log-Likelihood:                -53.677
Date:                Fri, 13 Nov 2015   Deviance:                       107.35
Time:                        11:24:05   Pearson chi2:                     102.
No. Iterations:                     6                                         
==============================================================================
                 coef    std err          z      P>|z|      [95.0% Conf. Int.]
------------------------------------------------------------------------------
const         -5.3095      1.134     -4.683      0.000        -7.531    -3.088
AGE            0.1109      0.024      4.610      0.000         0.064     0.158
==============================================================================
'''
glm_results.predict([[1,20],[1,23],[1,24]])

glm_results.mu[:3]



#%% Using statsmodels.formula.api.glm


dat=data.copy()
dta = dat[['CHD','AGE']]
endog = dta['CHD']
formula = 'CHD ~ AGE'

glm_lr = smf.glm(formula=formula, data=dta, family=sm.families.Binomial())
res = glm_lr.fit()
print(res.summary())
 
'''
 Generalized Linear Model Regression Results                  
==============================================================================
Dep. Variable:                    CHD   No. Observations:                  100
Model:                            GLM   Df Residuals:                       98
Model Family:                Binomial   Df Model:                            1
Link Function:                  logit   Scale:                             1.0
Method:                          IRLS   Log-Likelihood:                -53.677
Date:                Fri, 13 Nov 2015   Deviance:                       107.35
Time:                        14:00:02   Pearson chi2:                     102.
No. Iterations:                     6                                         
==============================================================================
                 coef    std err          z      P>|z|      [95.0% Conf. Int.]
------------------------------------------------------------------------------
Intercept     -5.3095      1.134     -4.683      0.000        -7.531    -3.088
AGE            0.1109      0.024      4.610      0.000         0.064     0.158
==============================================================================

'''
#%% Using scikit-learn
from sklearn import linear_model

dat=data.copy()
dta = dat[['CHD','AGE']]

model = linear_model.LogisticRegression(fit_intercept=False, penalty='l2', dual=False, tol=0.0001, C=10000)
dat['const'] = 1
x=dat[['const','AGE']]
y=dat['CHD']
model.fit(x,y)
model.coef_






# -*- coding: utf-8 -*-
"""
Created on Fri Nov 13 15:25:23 2015

@author: ub71894
"""
from pandas import Series, DataFrame 
import pandas as pd
import numpy as np


df1 = pd.DataFrame({'key': ['b', 'b', 'a', 'c', 'a', 'a', 'b'],'data1': range(7)})
df2 = pd.DataFrame({'key': ['a', 'b', 'd'],'data2': range(3)})
pd.merge(df1, df2)


left = pd.DataFrame({'key1': ['foo', 'foo', 'bar'], 'key2': ['one', 'two', 'one'],'lval': [1, 2, 3]})

right = pd.DataFrame({'key1': ['foo', 'foo', 'bar', 'bar'],'key2': ['one', 'one', 'one', 'two'],'rval': [4, 5, 6, 7]})


pd.merge(left, right, on=['key1', 'key2'], how='outer')


arr = np.arange(12).reshape((3, 4))

pd.concat([arr, arr])

s1 = pd.Series([0, 1], index=['a', 'b'])
s2 = pd.Series([2, 3, 4], index=['c', 'd', 'e'])
s3 = pd.Series([5, 6], index=['f', 'g'])

s4 = pd.concat([s1 * 5, s3])



data = DataFrame({'food': ['bacon', 'pulled pork', 'bacon', 'Pastrami',
'corned beef', 'Bacon', 'pastrami', 'honey ham',
'nova lox'],
'ounces': [4, 3, 12, 6, 7.5, 8, 3, 5, 6]})

meat_to_animal = {
'bacon': 'pig',
'pulled pork': 'pig',
'pastrami': 'cow',
'corned beef': 'cow',
'honey ham': 'pig',
'nova lox': 'salmon'
}
data['animal'] = data['food'].map(str.lower).map(meat_to_animal)

#%%

ages = [20, 22, 25, 27, 21, 23, 37, 31, 61, 45, 41, 32]
bins = [18, 25, 35, 60, 100]
cats = pd.cut(ages, bins)
group_names = ['Youth', 'YoungAdult', 'MiddleAged', 'Senior']
cats = pd.cut(ages, bins, labels=group_names)


data=np.random.randn(1000)
cc =pd.qcut(data,[0,0.1,0.5,0.9,1.0])
pd.value_counts(cc)
#%%

np.random.seed(12345)
data = DataFrame(np.random.randn(1000, 4))

data.describe()



df = DataFrame({'key': ['b', 'b', 'a', 'c', 'a', 'b'], 'data1': range(6)})
dummies = pd.get_dummies(df['key'], prefix='key')


#%%
import json
from pandas import Series, DataFrame 
import pandas as pd

db = json.load(open('foods-2011-10-03.json'))
nutrients = []

info_keys = ['description', 'group', 'id', 'manufacturer']
info = DataFrame(db, columns=info_keys)
col_mapping = {'description' : 'food', 'group' : 'fgroup'}
info = info.rename(columns=col_mapping, copy=False)

for rec in db:
    fnuts = DataFrame(rec['nutrients'])
    fnuts['id'] = rec['id']
    nutrients.append(fnuts)


nutrients = pd.concat(nutrients, ignore_index=True)
nutrients = nutrients.drop_duplicates()
col_mapping = {'description' : 'nutrient','group' : 'nutgroup'}
nutrients = nutrients.rename(columns=col_mapping, copy=False)

ndata = pd.merge(nutrients, info, on='id', how='outer')



result = ndata.groupby(['nutrient', 'fgroup'])['value'].quantile(0.5)



result['Zinc, Zn'].order().plot(kind='barh')

by_nutrient = ndata.groupby(['nutgroup', 'nutrient'])

get_maximum(c[['food','value']])
get_maximum(c[['value','food']])
get_maximum(c)
b[['food','value']]





tt = pd.DataFrame({'name': ['Yu', 'Zheng'],'age': [31, 32],'height':[173,164]})

tt2 = pd.DataFrame({'name': ['Yu2', 'Zheng'],'weight': [164, 110],'size':[10,8]})

tt3 = pd.DataFrame({'name': ['Huang', 'Lu'],'age': [50, 60],'height':[133,124]})

pd.merge(tt, tt2, how='outer',on='name')
pd.merge(tt, tt3,how='outer')











from sklearn.metrics import roc_auc_score


from sklearn.datasets import make_hastie_10_2
from sklearn.ensemble import GradientBoostingClassifier
from sklearn.ensemble.partial_dependence import plot_partial_dependence

#%% easy emample
X, y = make_hastie_10_2(random_state=0)
X_train, X_test = X[:2000], X[2000:]
y_train, y_test = y[:2000], y[2000:]

clf = GradientBoostingClassifier(n_estimators=100, learning_rate=1.0, max_depth=1, random_state=0).fit(X_train, y_train)
clf.score(X_train, y_train)               

yest = clf.predict(X_train)
roc_auc_score(y_train,yest)



#%% CRE PERM data
X_train = combo.iloc[:,:9]
y_train = combo.iloc[:,9]
clf = GradientBoostingClassifier(n_estimators=200, learning_rate=0.1, max_depth=10, random_state=0,subsample=0.5).fit(X_train, y_train)

yest = clf.predict(X_train)
roc_auc_score(y_train,yest)

features = [0,1,2,3,(1,3)]
fig, axs = plot_partial_dependence(clf, X_train, features, feature_names=('Eff_Gross_Income','LTV','Tot_cur_UBOC_Comm_over_NOI','Sbmkt_Vacancy_Rate_Pct')) 




#%%

from sklearn.ensemble import GradientBoostingClassifier


clf = GradientBoostingClassifier(n_estimators=100, learning_rate=0.1,max_depth=10, random_state=0).fit(X_train, y_train)
features = [1,3,(0,1)]
fig, axs = plot_partial_dependence(clf, X_train, features, feature_names=('Eff_Gross_Income', 'LTV')) 


#%%



















#%%
name=['AR_def_in','AR_PDRR_in','within2_PDRR_in','within1_PDRR_in','AR_PDRR_out','within2_PDRR_out','within1_PDRR_out,']
result = [AR_def_in,AR_PDRR_in,within2_PDRR_in,within1_PDRR_in,AR_PDRR_out,within2_PDRR_out,within1_PDRR_out]
df = pd.DataFrame.from_dict(dict(zip(name,result)))
df=df[name]import numba
import numpy as np
import time

@numba.jit
def sum2d_nb(arr):
    M, N = arr.shape
    result = 0.0
    for i in range(M):
        for j in range(N):
            result += arr[i,j]
    return result



def sum2d(arr):
    M, N = arr.shape
    result = 0.0
    for i in range(M):
        for j in range(N):
            result += arr[i,j]
    return result


a = np.random.rand(99999999).reshape(11111111,9)



start_time = time.time()
print(sum2d(a))
print("--- %s seconds ---" % (time.time() - start_time))


start_time = time.time()
print(sum2d_nb(a))
print("--- %s seconds ---" % (time.time() - start_time))
"""

"""
import os, numpy as np, pandas as pd
os.chdir(r"C:\Users\ub71894\Documents\code\python\testcode")

from PIL import Image
import pytesseract

pic = Image.open('captcha1.jpg')
print(pytesseract.image_to_string(pic))


print(pytesseract.image_to_string(Image.open('captcha1.png'), lang='fra'))'''
In this function, you have:
-   Turning format into percent
-   Changing police (bold, size…)
-   Defining vertical/horizontal alignment
-   Coloring
-   Row/ciolumn dimensions
-   Adding borders
'''

import os, numpy as np, pandas as pd
os.chdir(r"C:\Users\ub71894\Documents\code\python\testcode")

#%% use openpyxl to create new xlsx:

from openpyxl import Workbook
from openpyxl.compat import range
from openpyxl.cell import get_column_letter

wb = Workbook()
dest_filename = 'empty_book.xlsx'

ws1 = wb.active
ws1.title = "range names"

for row in range(1, 40):
    ws1.append(range(600))

ws2 = wb.create_sheet(title="Pi")
ws2['F5'] = 3.14

ws3 = wb.create_sheet(title="Data")
for row in range(10, 20):
    for col in range(27, 54):
        _ = ws3.cell(column=col, row=row, value="%s" % get_column_letter(col))

wb.save(filename = dest_filename)







#%% read specific range frpm existing file:

bf_existing = pd.read_excel('top5_AR_5%.xlsx',0,skiprows=1, parse_cols = 'B:H')
af_existing = pd.read_excel('top5_AR_5%.xlsx',0,skiprows=1, parse_cols = 'K:P')
af_new = pd.read_excel('top5_AR_5%.xlsx',0,skiprows=1, parse_cols = 'T:Y')























#%% condition formatting:

import numpy as np
import pandas
import openpyxl
from openpyxl import Workbook
from openpyxl.styles import Border, Side, Alignment, Font
from openpyxl.formatting.rule import IconSetRule
from openpyxl.styles import PatternFill

        
#Open output
wb =  openpyxl.load_workbook(output_file_path)
sheet = wb.create_sheet() 
sheet.title = u'RLA-Override'
    sheet.cell(row=3,column=col).value=RLA


def CreateStyleSheet_RLA_Override(sheet):
    
    for col in range(3,7):
        sheet.cell(row=3,column=col).number_format='.0%'
        sheet.cell(row=6,column=col).number_format='.0%'
        sheet.cell(row=4,column=col).number_format='.0%'    
        sheet.cell(row=5,column=col).number_format='.0%'
        
    #Work on style
    for _row in sheet.iter_rows('B1:F7'):
        for _cell in _row:
            _cell.alignment=Alignment(horizontal='center',vertical='center')
            _cell.font=Font(size=9)
            fill=PatternFill(start_color='FFFFFF', end_color='FFFFFF',fill_type='solid')
            _cell.fill=fill 
    for _row in sheet.iter_rows('C3:F3'): #Adjust color for cells
        for _cell in _row:
            color=Color(_cell.value, 'RLA')
            fill=PatternFill(start_color=color, end_color=color,fill_type='solid')
            _cell.fill=fill 
    for _row in sheet.iter_rows('C4:F4'): #Adjust color for cells
        for _cell in _row:
            color=Color(_cell.value, 'Override')
            fill=PatternFill(start_color=color, end_color=color,fill_type='solid')
            _cell.fill=fill
            
                    
    sheet.row_dimensions[2].height=17
    sheet.row_dimensions[3].height=17
    sheet.row_dimensions[4].height=17   
    sheet.row_dimensions[5].height=17    
    sheet.row_dimensions[6].height=17    
    sheet.row_dimensions[7].height=17
    sheet.column_dimensions['B'].width=14
    sheet.column_dimensions['C'].width=9
    sheet.column_dimensions['D'].width=9
    sheet.column_dimensions['E'].width=9   
    sheet.column_dimensions['F'].width=9
    
    for i in ['B1', 'B2','B3','B4', 'B5','B6', 'B7']:
        a = sheet[i]
        a.font=Font(bold=True, size=9)
        a.alignment=Alignment(vertical='center', horizontal='left')
    for i in ['B2', 'C2','D2','E2', 'F2']:
        a = sheet[i]
        a.font=Font(bold=True, size=9)
        a.alignment=Alignment(vertical='center', horizontal='center')
    sheet.merge_cells('B1:F1')
    sheet['B1'].alignment=Alignment(horizontal='center')
    
    #### Add borders ###
    side=Side(style='thin', color="FF000000")
    
    for i in ['C1','D1','E1']: #top border
        a = sheet[i]
        a.border=Border(bottom=side, top=side)
    sheet['F1'].border=Border(bottom=side, right=side, top=side)
    sheet['B1'].border=Border(bottom=side, left=side, top=side)
    
    
    for i in ['C2','D2','E2']: #top border
        a = sheet[i]
        a.border=Border(bottom=side, top=side)
    sheet['F2'].border=Border(bottom=side, right=side, top=side)
    sheet['B2'].border=Border(bottom=side, left=side, top=side, right=side)
    
    sheet['B3'].border=Border(left=side, top=side, right=side)
    sheet['F3'].border=Border(right=side, top=side)
    
    sheet['B4'].border=Border(left=side, right=side)
    sheet['F4'].border=Border(right=side)
    
    for i in ['C5','D5','E5']: #top border
        a = sheet[i]
        a.border=Border(top=side)
   sheet['B5'].border=Border(top=side, left=side, right=side)
    sheet['F5'].border=Border(right=side, top=side)    
    
    for i in ['C6','D6','E6']: #top border
        a = sheet[i]
        a.border=Border(bottom=side)
    sheet['B6'].border=Border(bottom=side, left=side, right=side)
    sheet['F6'].border=Border(bottom=side, right=side)      
    
    sheet['B7'].border=Border(bottom=side, right=side, left=side, top=side)      
    sheet['C7'].border=Border(bottom=side, left=side, top=side) 
    sheet['D7'].border=Border(bottom=side, top=side)       
    sheet['E7'].border=Border(bottom=side, top=side)  
    sheet['F7'].border=Border(bottom=side, top=side, right=side)
#%%
import numpy as np
from pandas import Series, DataFrame
import pandas as pd

obj=Series(range(4), index=['d', 'a', 'b', 'c'])


#%%

frame = DataFrame(np.arange(8).reshape((2, 4)), index=['three', 'one'],
.....: columns=['d', 'a', 'b', 'c'])

frame = DataFrame({'b': [4.3, 7, -3, 2], 'a': [0, 1, 0, 1],'c': [-2, 5, 8, -2.5]})


df = DataFrame(np.random.randn(4, 3), index=['a', 'a', 'b', 'b'])

df = DataFrame([[1.4, np.nan], [7.1, -4.5],
.....: [np.nan, np.nan], [0.75, -1.3]],
.....: index=['a', 'b', 'c', 'd'],
.....: columns=['one', 'two'])

#%%
dat = DataFrame(np.random.randn(100, 2), columns=['a', 'b'])

#%%

data = {'state':['Ohio','Ohio','Ohio','Nevada','Nevada'],'year':[2000,2001,2002,2001,2002],'pop':[1.5,1.7,3.6,2.4,2.9]}
df = DataFrame(data)

aa.groupby('year')['pop'].mean()

import os
dir_name=r'C:\Users\ub71894\Documents\code\Python\testcode'
os.chdir(dir_name)


from pptx import Presentation
from pptx.util import Inches, Pt

prs = Presentation()
blank_slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(blank_slide_layout)

txBox = slide.shapes.add_textbox(Inches(0.4), Inches(6.5), Inches(1), Inches(1))
tf = txBox.text_frame

p = tf.add_paragraph()
p.text = r"* Accuracy rate is represented by Somers'D"
p.font.bold = True
p.font.size = Pt(10)




prs.save('test.pptx')
# -*- coding: utf-8 -*-
"""
Created on Fri Mar 25 17:02:33 2016

@author: 
"""



import pyautogui,time

#pyautogui.PAUSE = 1
pyautogui.FAILSAFE = False

#pyautogui.size()
#width, height = pyautogui.size()
#%% office
i=0
while True:
    pyautogui.moveTo(1630, 547, duration=5)
    pyautogui.moveTo(1632, 547, duration=5)
    time.sleep(50)
    i+=1    
    print(i)

#%%no dock:
i=0
while True:
    time.sleep(57)
    pyautogui.moveTo(1778, 413, duration=3)
    pyautogui.click(1779,404)
    i+=1
    print(i)

#%%
    
i=0
while i<600:
    time.sleep(57)
    pyautogui.moveTo(1778, 413, duration=3)
    pyautogui.click(1779,404)
    i+=1
    print(i)
    
 #%%
pyautogui.position()

#%% timer to shutdown:
import os
i=0
while i<20:
    time.sleep(57)
    pyautogui.moveTo(1778, 413, duration=3)
    pyautogui.click(1779,404)
    i+=1
    print(i)
    
os.system("shutdown /s /t 1") 
import pymc3 as pm
import matplotlib.pyplot as plt
plt.style.use('seaborn')
import numpy as np
import pandas as pd

def logistic(x, beta, alpha=0):
    return 1.0 / (1.0 + np.exp(np.dot(beta, x) + alpha))

challenger_data = np.array([[ 66.,   0.],
 [ 70.,   1.],
 [ 69.,   0.],
 [ 68.,   0.],
 [ 67.,   0.],
 [ 72.,   0.],
 [ 73.,   0.],
 [ 70.,   0.],
 [ 57.,   1.],
 [ 63.,   1.],
 [ 70.,   1.],
 [ 78.,   0.],
 [ 67.,   0.],
 [ 53.,   1.],
 [ 67.,   0.],
 [ 75.,   0.],
 [ 70.,   0.],
 [ 81.,   0.],
 [ 76.,   0.],
 [ 79.,   0.],
 [ 75.,   1.],
 [ 76.,   0.],
 [ 58.,   1.]])


x = np.linspace(-4, 4, 100)

plt.plot(x, logistic(x, 1), label=r"$\beta_x^3 = 1$", ls="--", lw=1)
plt.plot(x, logistic(x, 3), label=r"$\beta = 3$", ls="--", lw=1)
plt.plot(x, logistic(x, -5), label=r"$\beta = -5$", ls="--", lw=1)

plt.plot(x, logistic(x, 1, 1), label=r"$\beta = 1, \alpha = 1$",
         color="#348ABD")
plt.plot(x, logistic(x, 3, -2), label=r"$\beta = 3, \alpha = -2$",
         color="#A60628")
plt.plot(x, logistic(x, -5, 7), label=r"$\beta = -5, \alpha = 7$",
         color="#7A68A6")

plt.legend(loc="lower left");


#%%
temperature = challenger_data[:, 0]
D = challenger_data[:, 1]  # defect or not?

#notice the`value` here. We explain why below.
with pm.Model() as model:
    beta = pm.Normal("beta", mu=0, tau=0.001, testval=0)
    alpha = pm.Normal("alpha", mu=0, tau=0.001, testval=0)
    p = pm.Deterministic("p", 1.0/(1. + tt.exp(beta*temperature + alpha)))

with model:
    observed = pm.Bernoulli("bernoulli_obs", p, observed=D)
    
    # Mysterious code to be explained in Chapter 3
    start = pm.find_MAP()
    step = pm.Metropolis()
    trace = pm.sample(12000, step=step, start=start)
    burned_trace = trace[10000::2]

#%%
alpha_samples = burned_trace["alpha"][:, None]  # best to make them 1d
beta_samples = burned_trace["beta"][:, None]
#histogram of the samples:
plt.subplot(211)
plt.title(r"Posterior distributions of the variables $\alpha, \beta$")
plt.hist(beta_samples, histtype='stepfilled', bins=35, alpha=0.85,
         label=r"posterior of $\beta$", color="#7A68A6", normed=True)
plt.legend()

plt.subplot(212)
plt.hist(alpha_samples, histtype='stepfilled', bins=35, alpha=0.85,
         label=r"posterior of $\alpha$", color="#A60628", normed=True)
plt.legend();


#%%
t = np.linspace(temperature.min() - 5, temperature.max()+5, 50)[:, None]
p_t = logistic(t.T, beta_samples, alpha_samples)
mean_prob_t = p_t.mean(axis=0)

plt.plot(t, mean_prob_t, lw=3, label="average posterior \nprobability \
of defect")
plt.plot(t, p_t[136, :], ls="--", label="realization from posterior")
plt.plot(t, p_t[1943, :], ls="--", label="realization from posterior")
plt.scatter(temperature, D, color="k", s=50, alpha=0.5)
plt.title("Posterior expected value of probability of defect; \
plus realizations")
plt.legend(loc="lower left")
plt.ylim(-0.1, 1.1)
plt.xlim(t.min(), t.max())
plt.ylabel("probability")
plt.xlabel("temperature");


#%%
from scipy.stats.mstats import mquantiles

# vectorized bottom and top 2.5% quantiles for "confidence interval"
qs = mquantiles(p_t, [0.025, 0.975], axis=0)
plt.fill_between(t[:, 0], *qs, alpha=0.7,
                 color="#7A68A6")

plt.plot(t[:, 0], qs[0], label="95% CI", color="#7A68A6", alpha=0.7)

plt.plot(t, mean_prob_t, lw=1, ls="--", color="k",
         label="average posterior \nprobability of defect")

plt.xlim(t.min(), t.max())
plt.ylim(-0.02, 1.02)
plt.legend(loc="lower left")
plt.scatter(temperature, D, color="k", s=50, alpha=0.5)
plt.xlabel("temp, $t$")

plt.ylabel("probability estimate")
plt.title("Posterior probability estimates given temp. $t$");


#%%

N = 10000
with pm.Model() as model:
    beta = pm.Normal("beta", mu=0, tau=0.001, testval=0)
    alpha = pm.Normal("alpha", mu=0, tau=0.001, testval=0)
    p = pm.Deterministic("p", 1.0/(1. + tt.exp(beta*temperature + alpha)))
    observed = pm.Bernoulli("bernoulli_obs", p, observed=D)
    
    simulated = pm.Bernoulli("bernoulli_sim", p, shape=p.tag.test_value.shape)
    step = pm.Metropolis(vars=[p])
    trace = pm.sample(N, step=step)

simulations = trace["bernoulli_sim"]
print(simulations.shape)

plt.title("Simulated dataset using posterior parameters")

for i in range(4):
    ax = plt.subplot(4, 1, i+1)
    plt.scatter(temperature, simulations[1000*i, :], color="k",
                s=50, alpha=0.6)


#%% separation plots
posterior_probability = simulations.mean(axis=0)
print("posterior prob of defect | realized defect ")
for i in range(len(D)):
    print("%.2f                     |   %d" % (posterior_probability[i], D[i]))

ix = np.argsort(posterior_probability)
print("probb | defect ")
for i in range(len(D)):
    print("%.2f  |   %d" % (posterior_probability[ix[i]], D[ix[i]]))


def separation_plot( p, y, **kwargs ):
    """
    This function creates a separation plot for logistic and probit classification. 
    See http://mdwardlab.com/sites/default/files/GreenhillWardSacks.pdf
    
    p: The proportions/probabilities, can be a nxM matrix which represents M models.
    y: the 0-1 response variables.
    
    """    
    assert p.shape[0] == y.shape[0], "p.shape[0] != y.shape[0]"
    n = p.shape[0]

    try:
        M = p.shape[1]
    except:
        p = p.reshape( n, 1 )
        M = p.shape[1]

    colors_bmh = np.array( ["#eeeeee", "#348ABD"] )


    fig = plt.figure( )
    
    for i in range(M):
        ax = fig.add_subplot(M, 1, i+1)
        ix = np.argsort( p[:,i] )
        #plot the different bars
        bars = ax.bar( np.arange(n), np.ones(n), width=1.,
                color = colors_bmh[ y[ix].astype(int) ], 
                edgecolor = 'none')
        ax.plot( np.arange(n+1), np.append(p[ix,i], p[ix,i][-1]), "k",
                 linewidth = 1.,drawstyle="steps-post" )
        #create expected value bar.
        ax.vlines( [(1-p[ix,i]).sum()], [0], [1] )
        plt.xlim( 0, n)
        
    plt.tight_layout()
    
    return


#%% chap 3
N = 1

# the true parameters, but of course we do not see these values...
lambda_1_true = 1
lambda_2_true = 3

#...we see the data generated, dependent on the above two values.
data = np.concatenate([
    stats.poisson.rvs(lambda_1_true, size=(N, 1)),
    stats.poisson.rvs(lambda_2_true, size=(N, 1))
], axis=1)
print("observed (2-dimensional,sample size = %d):" % N, data)

# plotting details.
x = y = np.linspace(.01, 5, 100)
likelihood_x = np.array([stats.poisson.pmf(data[:, 0], _x)
                        for _x in x]).prod(axis=1)
likelihood_y = np.array([stats.poisson.pmf(data[:, 1], _y)
                        for _y in y]).prod(axis=1)
L = np.dot(likelihood_x[:, None], likelihood_y[None, :])


plt.subplot(221)
uni_x = stats.uniform.pdf(x, loc=0, scale=5)
uni_y = stats.uniform.pdf(x, loc=0, scale=5)
M = np.dot(uni_x[:, None], uni_y[None, :])
im = plt.imshow(M, interpolation='none', origin='lower',
                cmap=jet, vmax=1, vmin=-.15, extent=(0, 5, 0, 5))
plt.scatter(lambda_2_true, lambda_1_true, c="k", s=50, edgecolor="none")
plt.xlim(0, 5)
plt.ylim(0, 5)
plt.title("Landscape formed by Uniform priors on $p_1, p_2$.")

plt.subplot(223)
plt.contour(x, y, M * L)
im = plt.imshow(M * L, interpolation='none', origin='lower',
                cmap=jet, extent=(0, 5, 0, 5))
plt.title("Landscape warped by %d data observation;\n Uniform priors on $p_1, p_2$." % N)
plt.scatter(lambda_2_true, lambda_1_true, c="k", s=50, edgecolor="none")
plt.xlim(0, 5)
plt.ylim(0, 5)

plt.subplot(222)
exp_x = stats.expon.pdf(x, loc=0, scale=3)
exp_y = stats.expon.pdf(x, loc=0, scale=10)
M = np.dot(exp_x[:, None], exp_y[None, :])

plt.contour(x, y, M)
im = plt.imshow(M, interpolation='none', origin='lower',
                cmap=jet, extent=(0, 5, 0, 5))
plt.scatter(lambda_2_true, lambda_1_true, c="k", s=50, edgecolor="none")
plt.xlim(0, 5)
plt.ylim(0, 5)
plt.title("Landscape formed by Exponential priors on $p_1, p_2$.")

plt.subplot(224)
# This is the likelihood times prior, that results in the posterior.
plt.contour(x, y, M * L)
im = plt.imshow(M * L, interpolation='none', origin='lower',
                cmap=jet, extent=(0, 5, 0, 5))

plt.scatter(lambda_2_true, lambda_1_true, c="k", s=50, edgecolor="none")
plt.title("Landscape warped by %d data observation;\n Exponential priors on \
$p_1, p_2$." % N)
plt.xlim(0, 5)
plt.ylim(0, 5);


#%%
data = np.loadtxt(r'C:\Users\ub71894\Documents\Code\Python\Testcode\BMH\Chapter3_MCMC\data\mixture_data.csv', delimiter=",")

plt.hist(data, bins=20, color="k", histtype="stepfilled", alpha=0.8)
plt.title("Histogram of the dataset")
plt.ylim([0, None]);
print(data[:10], "...")


import theano.tensor as T


with pm.Model() as model:
    p1 = pm.Uniform('p', 0, 1)
    p2 = 1 - p1
    p = T.stack([p1, p2])
    assignment = pm.Categorical("assignment", p, 
                                shape=data.shape[0],
                                testval=np.random.randint(0, 2, data.shape[0]))
    
print("prior assignment, with p = %.2f:" % p1.tag.test_value)
print(assignment.tag.test_value[:10])


with model:
    sds = pm.Uniform("sds", 0, 100, shape=2)
    centers = pm.Normal("centers", 
                        mu=np.array([120, 190]), 
                        sd=np.array([10, 10]), 
                        shape=2)
    
    center_i = pm.Deterministic('center_i', centers[assignment])
    sd_i = pm.Deterministic('sd_i', sds[assignment])
    
    # and to combine it with the observations:
    observations = pm.Normal("obs", mu=center_i, sd=sd_i, observed=data)
    
print("Random assignments: ", assignment.tag.test_value[:4], "...")
print("Assigned center: ", center_i.tag.test_value[:4], "...")
print("Assigned standard deviation: ", sd_i.tag.test_value[:4])

with model:
    step1 = pm.Metropolis(vars=[p, sds, centers])
    step2 = pm.ElemwiseCategorical(vars=[assignment])
    trace = pm.sample(25000, step=[step1, step2])


































import os, sys, pandas as pd, numpy as np
os.chdir(r"C:\Users\ub71894\Documents\Docs\cheatsheet")
from PyPDF2 import PdfFileMerger

pdfs = ['Bokeh.pdf',
'ImportingData.pdf',
'JupyterNotebook.pdf',
'Keras.pdf',
'Matplotlib.pdf',
'Numpy.pdf',
'Pandas.pdf',
'PySpark_RDD.pdf',
'PySpark_SQL.pdf',
'ScikitLearn.pdf',
'SciPy_LinearAlgebra.pdf',
'Seaborn.pdf']


merger = PdfFileMerger()

for pdf in pdfs:
    merger.append(pdf)

merger.write("result.pdf")import pyupset as pyu

import os, sys, pandas as pd, numpy as np
os.chdir(r"C:\Users\ub71894\Documents\code\python\Testcode")


set1 = [1,2,3,4,5,6]
set2 = [2,3,4,7,8,9]
set3 = [1,3,5,7,9,10]
df1=pd.DataFrame()
df2=pd.DataFrame()
df3=pd.DataFrame()
df1['num'] = set1
df2['num'] = set2
df3['num'] = set3

df = {'set1':df1, 'set2':df2, 'set3':df3}
pyu.plot(df)


with open('test_data_dict.pckl', 'rb') as f:
    data_dict = load(f)
pyu.plot(data_dict)

#%%

from sas7bdat import SAS7BDAT
import pandas as pd 
import os

os.chdir(r"C:\Users\ub71894\Documents\code\python\testcode")


#old mehtod
f = SAS7BDAT('cola.sas7bdat') 
ff= f.to_data_frame()


#new: doesn't work

df = pd.read_sas("cola.sas7bdat") 
# -*- coding: utf-8 -*-
"""
Created on Fri Jan 29 15:54:08 2016

@author: ub71894 (4e8e6d0b), CSG
"""


import os
import pandas as pd 
import numpy as np
os.chdir(r"C:\Users\ub71894\Documents\code\Python\testcode")


# make up 2 DataFrame
dat = pd.read_csv(r"C:\Users\ub71894\Documents\data\RA\ra_perm_csv.csv")
a=dat.iloc[1:100,1:5].copy()
b=dat.iloc[100:200,6:10].copy()
c=a


#%% save df to an excel file:
writer = pd.ExcelWriter('temp.xlsx', engine='xlsxwriter')
a.to_excel(writer, sheet_name='DataFrame_a')
b.to_excel(writer, sheet_name='DataFrame_b')
writer.save()


#%% save df to an existing excel file:

from openpyxl import load_workbook

book = load_workbook('temp.xlsx')
writer = pd.ExcelWriter('temp.xlsx', engine='openpyxl') 
writer.book = book
writer.sheets = dict((ws.title, ws) for ws in book.worksheets)

c.to_excel(writer, sheet_name='DataFrame_c')

writer.save()



#%% write number in specific range of workbook:
from openpyxl import Workbook
from openpyxl.compat import range

wb = Workbook()

dest_filename = 'empty_book.xlsx'

ws1 = wb.active
ws1.title = "range names"

for row in range(1, 40):
    ws1.append(range(600))

ws2 = wb.create_sheet(title="Pi")

ws2['F5'] = df2

ws3 = wb.create_sheet(title="Data")
for row in range(10, 20):
    for col in range(27, 54):
        _ = ws3.cell(column=col, row=row, value="%s" % get_column_letter(col))
print(ws3['AA10'].value)

wb.save(filename = dest_filename)


#%% write df in specific range in each setting:
top = 5
writer = pd.ExcelWriter('top5.xlsx')
sheetname = ['dev','prod','combo']
ab= ['/before','/after']
filename = ['output\mfa_before.xlsx','output\mfa_after_existing.xlsx']

for i,file in enumerate(filename):
    for j,sheet in enumerate(sheetname):
        df = pd.read_excel(file,sheetname=sheet)
        df.sort_values(by='AR',ascending=False, inplace=True)
        df2 = df.iloc[:top,:]
        df2.to_excel(writer, sheet_name='Sheet1', startrow=(8*j+1), startcol=(10**i))
        worksheet = writer.sheets['Sheet1']
        worksheet.write((8*j+1), (10**i), sheet+ab[i])

writer.save()from twilio.rest import TwilioRestClient 
 
# put your own credentials here 
ACCOUNT_SID = "AC0c8290c2ea6e80979b9d419029e25c4f" 
AUTH_TOKEN = "9f45c5b4dd1428470aaa4e7ae5fa4707" 
 
client = TwilioRestClient(ACCOUNT_SID, AUTH_TOKEN) 
 
client.messages.create(
    to="16318355443", 
    from_="+16319542023", 
    body="测试",  
    #media_url=r'http://www.clker.com/cliparts/d/b/2/7/13440412741381521061sunflower1-md.png'
)



import smtplib
smtpObj = smtplib.SMTP('smtp.gmail.com', 587)
type(smtpObj)

smtpObj.ehlo()
smtpObj.starttls()

MY_SECRET_PASSWORD = input('Enter your ps: ')
smtpObj.login('wwhome16@gmail.com', MY_SECRET_PASSWORD)



smtpObj.sendmail('wwhome16@gmail.com', 'yuyangfirst@gmail.com',
'Subject: So long.\nDear Alice, so long and thanks for all the fish. Sincerely,dsasd \
sdfsdfsdfs fsdfsdfsd')
{}



smtpObj.quit()import os
import pandas as pd 
import numpy as np
os.chdir(r"C:\Users\ub71894\Documents\code\Python\testcode")

import win32com.client
myPowerPoint = win32com.client.DispatchEx( 'Powerpoint.Application' )

ppt = r'C:\Users\ub71894\Documents\code\Python\testcode\OriginalTemplate.pptx'
thePresentation = myPowerPoint.Presentations.Open(ppt, False, False, False )

for i in range(2):
    thePresentation.Slides(2).Duplicate()

thePresentation.SaveAs(r'C:\Users\ub71894\Documents\code\Python\testcode\OriginalTemplate_dup.pptx')
thePresentation.Close() 


"""
Minimal Example
===============
Generating a square wordcloud from the US constitution using default arguments.
"""

from os import path
from wordcloud import WordCloud

d = path.dirname(__file__)

# Read the whole text.
text = open(path.join(d, 'constitution.txt')).read()

# Generate a word cloud image
wordcloud = WordCloud().generate(text)

# Display the generated image:
# the matplotlib way:
import matplotlib.pyplot as plt
plt.imshow(wordcloud)
plt.axis("off")

# take relative word frequencies into account, lower max_font_size
wordcloud = WordCloud(max_font_size=40, relative_scaling=.5).generate(text)
plt.figure()
plt.imshow(wordcloud)
plt.axis("off")
plt.show()

# The pil way (if you don't have matplotlib)
#image = wordcloud.to_image()
#image.show()['autogui_search.py', 'compustat_db_build.py', 'compustat_db_build_yang.py', 'databaseCreation_CnI.py', 'Date Process.py', 'group.py', 'plot.py', 'python_test.py', 'queryexample.py', 'queryexample_yang.py', 'readtxt.py', 'regression&plot_test.py', 'reset.py', 'run_sql.py', 'scipy_test.py', 'seaborn_test.py', 'shenlu.py', 'subsetsum.py', 'sync.py', 'test_AZ_ranking.py', 'test_class.py', 'test_datetime.py', 'test_GAM.py', 'test_graphviz.py', 'test_gredientdecent_logiticreg.py', 'test_LogisticReg.py', 'test_merge.py', 'test_ML.py', 'test_numba.py', 'test_OCR.py', 'test_openpyxl.py', 'test_pandas.py', 'test_pptx.py', 'test_pyautogui.py', 'test_pymc.py', 'test_pypdf2.py', 'test_pyupset.py', 'test_sas7bdat.py', 'test_toexcel.py', 'test_twilio.py', 'test_win32ppt.py', 'test_wordcloud.py']
[180, 234, 288, 710, 792, 804, 914, 1203, 1233, 1269, 1341, 1460, 1463, 1498, 1582, 1663, 1738, 1762, 1838, 1931, 1968, 1982, 2069, 2081, 2191, 2315, 2455, 2524, 2559, 2573, 2751, 2783, 2808, 2866, 3216, 3242, 3265, 3282, 3369, 3405, 3422, 3455]
