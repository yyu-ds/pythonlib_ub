
import os
dir_name=r'C:\Users\ub71894\Documents\code\Python\testcode'
os.chdir(dir_name)


from pptx import Presentation
from pptx.util import Inches, Pt

prs = Presentation()
blank_slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(blank_slide_layout)

txBox = slide.shapes.add_textbox(Inches(0.4), Inches(6.5), Inches(1), Inches(1))
tf = txBox.text_frame

p = tf.add_paragraph()
p.text = r"* Accuracy rate is represented by Somers'D"
p.font.bold = True
p.font.size = Pt(10)




prs.save('test.pptx')
