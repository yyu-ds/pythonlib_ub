# -*- coding: utf-8 -*-
"""

Created on Wed Mar  2 12:29:17 2016

@author: xu79799 
@author: ub71894, All work after Jun,2016

"""

import pandas as pd
import numpy as np
import openpyxl
from openpyxl.styles import PatternFill
from openpyxl.utils import get_column_letter
from openpyxl.styles import Border, Side, Alignment, Font
from openpyxl.formatting.rule import IconSetRule
from openpyxl import load_workbook
import pandas.tseries.offsets as pto
from openpyxl.utils.cell import coordinate_from_string, column_index_from_string


def range_to_index(range_in_str): 

    delimiter = range_in_str.find(':')
    min_cell = range_in_str[:delimiter]
    max_cell = range_in_str[(delimiter+1):]

    xy = coordinate_from_string(min_cell) 
    min_col = column_index_from_string(xy[0]) 
    min_row = xy[1]
    xy = coordinate_from_string(max_cell) 
    max_col = column_index_from_string(xy[0]) 
    max_row = xy[1]

    return({'min_col':min_col, 'max_col':max_col, 'min_row':min_row, 'max_row':max_row})



## Fonction which take a rating (ex Aaa) and return a corresponding number 
def UBagencymap(Agency,WhichAgency, BusinessPDRRDict):
    Agency_loc = int(BusinessPDRRDict[Agency])
    return(Agency_loc)
    

#Count the number of obligors which have same rating with a margin of y    
def within(MM,y):
    x=np.array(MM)
    N=len(x)
    sum1=0
    if y == 0:
        return np.sum(np.diag(x))
    else:
        for i in range(0,N):
            ll=max(0,i-y)
            uu=min(N+1,i+y+1)
            sum1=sum1+np.sum(x[ll:uu,i])
        return(sum1)
        
def better_ratings(MM):
    x=np.array(MM)
    N=len(x)
    sum1=0
    for i in range(0,N):
        ll=0
        sum1=sum1+np.sum(x[ll:(i),i])
    return(sum1)
    
def worse_ratings(MM):
    x=np.array(MM)
    N=len(x)
    sum1=0
    for i in range(0,N):
        uu=N+1
        sum1=sum1+np.sum(x[(i+1):uu,i])
    return(sum1)
    
def GetLastRatings(df, obligor):
    res=df[df['CUSTOMERID']==obligor]
    
    n=res.shape[0]
    if n==0: raise 'Exposure asked for an obligor not included'
    last=res.iloc[[n-1]] 
    return last
    

def SelectDataForPrelimMatrix(df): #debug on 20161028  
    # modified on 20180424. to change 'prelim-final' TM to 'before JBA-final' TM
    df['RLA_Notches'].fillna(0, inplace=True)
    df['Override_Action'].fillna(0, inplace=True)
    df['JBA'] = df['RLA_Notches']+df['Override_Action']
    df['PDRR_Before_JBA'] = df['Final_PD_Risk_Rating'] - df['JBA']

    if 'CREPROPERTYID' in df.columns.tolist():
        cols = [ u'CUSTOMERID', 'CREPROPERTYID',u'PDRR_Before_JBA', u'Final_PD_Risk_Rating']
        cols_new = ['Obligor','Property','BeforeRating','AfterRating']
        to_dropdup = ['CUSTOMERID', 'CREPROPERTYID', 'archive_date']
    else:
        cols = [ u'CUSTOMERID', u'PDRR_Before_JBA', u'Final_PD_Risk_Rating']
        cols_new = ['Obligor','BeforeRating','AfterRating']
        to_dropdup = ['CUSTOMERID', 'archive_date']
        
    df.sort_values(by=to_dropdup, ascending=True, inplace=True)
    df.drop_duplicates(subset=to_dropdup, keep='last', inplace=True)
    df.dropna(subset=['Final_PD_Risk_Rating'], inplace=True)    

    recentsnapshot = pd.to_datetime(df.snapshot.max(), format='%Y%m') + pto.MonthEnd() # make it to the end of snapshot month
    df=df[ ((recentsnapshot- pto.DateOffset(years=1))< df[u'archive_date']) & \
    (df[u'archive_date']<= recentsnapshot)]

    # keep all columns  #20191022
    df.rename(columns=dict(zip(cols,cols_new)), inplace=True)
    
    return df




def SelectDataForMigration(df):
    if 'CREPROPERTYID' in df.columns.tolist():
        to_dropdup = ['CUSTOMERID', 'CREPROPERTYID', 'archive_date']
        to_merge = ['CUSTOMERID', 'CREPROPERTYID']
        cols = ['CUSTOMERID','CREPROPERTYID','Final_PD_Risk_Rating_y','Final_PD_Risk_Rating_x']
        cols_new = {'CUSTOMERID':'Obligor','CREPROPERTYID':'Property','Final_PD_Risk_Rating_y':'BeforeRating','Final_PD_Risk_Rating_x':'AfterRating'}
    else:
        to_dropdup = ['CUSTOMERID', 'archive_date']
        to_merge = 'CUSTOMERID'
        cols = ['CUSTOMERID','Final_PD_Risk_Rating_y','Final_PD_Risk_Rating_x']
        cols_new = {'CUSTOMERID':'Obligor','Final_PD_Risk_Rating_y':'BeforeRating','Final_PD_Risk_Rating_x':'AfterRating'}
    df.sort_values(by=to_dropdup, ascending=True, inplace=True)
    df.drop_duplicates(subset=to_dropdup, keep='last', inplace=True)
    df.dropna(subset=['Final_PD_Risk_Rating'], inplace=True)    


    recentsnapshot = pd.to_datetime(df.snapshot.max(), format='%Y%m') + pto.MonthEnd() # make it to the end of snapshot month
    currentyear=df[ ((recentsnapshot- pto.DateOffset(years=1))< df[u'archive_date']) & \
                    (df[u'archive_date']<= recentsnapshot)]
    previousyear = df[ ((recentsnapshot- pto.DateOffset(years=2))< df[u'archive_date']) & \
                    (df[u'archive_date']<= (recentsnapshot- pto.DateOffset(years=1)))]
    
    new = pd.merge(currentyear.drop_duplicates(subset=to_merge,keep='last'),\
        previousyear.drop_duplicates(subset=to_merge,keep='last'), \
        on=to_merge, how='inner')

    new.rename(columns=cols_new,inplace=True)  
    return new



def ComputeListOutstandingObligors(df):
    res=[]
    
    for i in range(df.shape[0]):
        if (abs(df['BeforeRating'].iloc[i]-df['AfterRating'].iloc[i])>=4):
            res.append(df['Obligor'].iloc[i])
    return res

def ComputeListOutstandingObligors_new(df):

    # created on 20191022 to apply new rule for outliers  
    df['diff'] = np.abs(df['BeforeRating'] - df['AfterRating'])
    df_part1 = df.query('AfterRating<=7 and diff>=4')
    df_part2 = df.query('8<=AfterRating<=12 and diff>=3')
    df_part3 = df.query('AfterRating>=13 and diff>=2')
    df_temp = pd.concat([df_part1, df_part2, df_part3])

    return df_temp


def FillExcelOutliers(df, outlier_file_path, df_outstanding, name_model, typeof):
    # modified on 20191022 to output all columns for outlier dataframe

    if 'CREPROPERTYID' in df.columns.tolist():
        Variables_to_be_printed=['CUSTOMERID','CREPROPERTYID','NET_OUTSTANDING_AMOUNT','EAD','Customer_Name','NET_COMMITMENT_TOT', 'archive_date','Prelim_PD_Risk_Rating','Final_PD_Risk_Rating', ]+['Override_Justification','Override_Reason_1','Override_Reason_2','Override_Reason_3','RLA_Justification','RLA_Reason_1','RLA_Reason_2','RLA_Reason_3']

    else:
        Variables_to_be_printed=['CUSTOMERID', 'NET_OUTSTANDING_AMOUNT','EAD','Customer_Name','NET_COMMITMENT_TOT', 'archive_date','Prelim_PD_Risk_Rating','Final_PD_Risk_Rating', ]+['Override_Justification','Override_Reason_1','Override_Reason_2','Override_Reason_3','RLA_Justification','RLA_Reason_1','RLA_Reason_2','RLA_Reason_3']
    
    if typeof=='Prelim':
        sheetname = name_model[:20]+'_BeforeJBA'
    else:
        sheetname = name_model[:20]+'_Migration'

    if len(df_outstanding)==0: 
        pass
    else:
        #df_outstanding.to_excel(writer_outlier, sheetname)


        wb =  openpyxl.load_workbook(outlier_file_path)
        sheet = wb.create_sheet() 
        sheet.title = sheetname
        
        book = load_workbook(outlier_file_path)
        writer = pd.ExcelWriter(outlier_file_path, engine='openpyxl')
        writer.book = book
        writer.sheets = dict((ws.title, ws) for ws in book.worksheets)
        df_outstanding.to_excel(writer, sheetname, index=False)
        writer.save()


def CreateStyleSheet_TransitionMatrix(sheet, rowshift, extra_c, colshift):
    
    PDRR= range(1,16)
    
    #### Work on style ###
    for i in range(rowshift-3,rowshift+15):
        sheet.row_dimensions[i].height=12
    for j in ['V','W','X']:
        sheet.column_dimensions[j].width=7
    sheet.column_dimensions['B'].width=3
    sheet.column_dimensions['U'].width=9
    sheet.column_dimensions['AA'].width=9
    sheet.column_dimensions['AB'].width=5
    sheet.column_dimensions['Y'].width=8
    sheet.column_dimensions['Z'].width=8
    

    for _row in sheet.iter_rows(**range_to_index('B'+ str(rowshift-2) +':R'+ str(rowshift-1))):
        for _cell in _row:
            fill=PatternFill(start_color='FFFFFF', end_color='FFFFFF',fill_type='solid')
            _cell.fill=fill 
    for _row in sheet.iter_rows(**range_to_index('B'+ str(rowshift-2)+':C'+ str(rowshift+15))):
        for _cell in _row:
            fill=PatternFill(start_color='FFFFFF', end_color='FFFFFF',fill_type='solid')
            _cell.fill=fill 
    
    for _row in sheet.iter_rows(**range_to_index('C' + str(rowshift-1)+ ':R'+ str(rowshift+15))):
        for _cell in _row:
            _cell.alignment=Alignment(horizontal='center',vertical='center')
            _cell.font=Font(size=9)
    
    for colNum in range(15):
         col1=colshift+colNum
         sheet.column_dimensions[get_column_letter(col1)].width = 5   
   
   #Format of percentage table
    for _row in sheet.iter_rows(**range_to_index('U'+ str(rowshift-1) +':AB'+ str(rowshift))): 
       for _cell in _row:
           _cell.font=Font(size=9, bold=True)
           _cell.alignment=Alignment(vertical='center', horizontal='center')
           fill=PatternFill(start_color='FFFFFF', end_color='FFFFFF',fill_type='solid')
           _cell.fill=fill 
           
    #Add borders:
    side=Side(style='thin', color="FF000000")
    for _row in sheet.iter_rows(**range_to_index('V'+ str(rowshift-1) +':X'+ str(rowshift-1))): #top border
       for _cell in _row:
           _cell.border=Border(top=side) 
    for _row in sheet.iter_rows(**range_to_index('V'+ str(rowshift) +':X'+ str(rowshift))): #bottom border
       for _cell in _row:
           _cell.border=Border(bottom=side)
    sheet['U'+str(rowshift-1)].border=Border(left=side, top=side)      
    sheet['U'+str(rowshift)].border=Border(left=side, bottom=side)
    sheet['Y'+str(rowshift-1)].border=Border(right=side, top=side)      
    sheet['Y'+str(rowshift)].border=Border(right=side, bottom=side)
    sheet['Z'+str(rowshift-1)].border=Border(left=side, top=side) 
    sheet['Z'+str(rowshift)].border=Border(left=side, bottom=side)
    sheet['AA'+str(rowshift-1)].border=Border(right=side, top=side) 
    sheet['AA'+str(rowshift)].border=Border(right=side, bottom=side)    
    sheet['AB'+str(rowshift-1)].border=Border(right=side, top=side, left=side) 
    sheet['AB'+str(rowshift)].border=Border(right=side, bottom=side, left=side)
    
### Do formatting ###   modified on 02/25/2019
    for rowNum in PDRR:
        rowNum=rowNum-1
        for colNum in range(len(PDRR)): 
            row1=rowshift+rowNum
            col1=colshift+colNum
            if rowNum==colNum:
                fill=PatternFill(start_color='FFFF00', end_color='FFFF00',fill_type='solid')
                sheet.cell(row=row1,column=col1).fill=fill
            elif (((abs(rowNum-colNum)== 1) | (abs(rowNum-colNum)== 2)) & (colNum<=6)):
                fill=PatternFill(start_color='FFFFCC', end_color='FFFFCC',fill_type='solid')
                sheet.cell(row=row1,column=col1).fill=fill
            elif ((abs(rowNum-colNum)== 3) & (colNum<=6)):
                fill=PatternFill(start_color='FCD5B4', end_color='FCD5B4',fill_type='solid')
                sheet.cell(row=row1,column=col1).fill=fill
            elif ((abs(rowNum-colNum)>3) & (not(sheet.cell(row=row1,column=col1).value==None)) & (colNum<=6)): 
                fill=PatternFill(start_color='E4DFEC', end_color='E4DFEC',fill_type='solid')
                sheet.cell(row=row1,column=col1).fill=fill 

            elif ((abs(rowNum-colNum)==1)  & (colNum<=11)):
                fill=PatternFill(start_color='FFFFCC', end_color='FFFFCC',fill_type='solid')
                sheet.cell(row=row1,column=col1).fill=fill
            elif ((abs(rowNum-colNum)== 2) & (colNum<=11)):
                fill=PatternFill(start_color='FCD5B4', end_color='FCD5B4',fill_type='solid')
                sheet.cell(row=row1,column=col1).fill=fill
            elif ((abs(rowNum-colNum)>2) & (not(sheet.cell(row=row1,column=col1).value==None)) & (colNum<=11)): 
                fill=PatternFill(start_color='E4DFEC', end_color='E4DFEC',fill_type='solid')
                sheet.cell(row=row1,column=col1).fill=fill 

            elif ((abs(rowNum-colNum)== 1) & (colNum>11)):
                fill=PatternFill(start_color='FCD5B4', end_color='FCD5B4',fill_type='solid')
                sheet.cell(row=row1,column=col1).fill=fill
            elif ((abs(rowNum-colNum)>1) & (not(sheet.cell(row=row1,column=col1).value==None)) & (colNum>11)): 
                fill=PatternFill(start_color='E4DFEC', end_color='E4DFEC',fill_type='solid')
                sheet.cell(row=row1,column=col1).fill=fill 

            else:
                fill=PatternFill(start_color='FFFFFF', end_color='FFFFFF',fill_type='solid')
                sheet.cell(row=row1,column=col1).fill=fill 
    
    if extra_c==0:
        
        #Title format        
        sheet.merge_cells('B'+str(rowshift-3)+':R'+str(rowshift-3))
        sheet.merge_cells('D'+str(rowshift-2)+ ':R'+str(rowshift-2))
        sheet.merge_cells('B'+str(rowshift)+ ':B'+str(rowshift+14))
        sheet['B'+str(rowshift)].alignment=Alignment(horizontal='center',text_rotation=90, vertical='center')
        sheet['B'+str(rowshift-3)].alignment=Alignment(horizontal='center', vertical='center')
        sheet['D'+str(rowshift-2)].alignment=Alignment(horizontal='center', vertical='center')
        sheet['B'+str(rowshift-3)].font=Font(bold=True)
        sheet['D'+str(rowshift-2)].font=Font(italic=True)
        sheet['B'+str(rowshift)].font=Font(italic=True)
        
        #Add borders
        for _row in sheet.iter_rows(**range_to_index('B'+ str(rowshift-3) +':R'+ str(rowshift-3))): #top border
           for _cell in _row:
               _cell.border=Border(top=side)
        for _row in sheet.iter_rows(**range_to_index('B'+ str(rowshift) +':R'+ str(rowshift))): #top border 2
           for _cell in _row:
               _cell.border=Border(top=side)
        for i in ['B'+ str(rowshift-2),'B'+ str(rowshift-1)]: #left border
            a = sheet[i]
            a.border=Border(left=side)
        for _row in sheet.iter_rows(**range_to_index('B'+ str(rowshift) +':B'+ str(rowshift+14))): #left border
           for _cell in _row:
               _cell.border=Border(left=side, right=side)
        for _row in sheet.iter_rows(**range_to_index('C'+ str(rowshift+14) +':R'+ str(rowshift+14))): #bottom border
           for _cell in _row:
               _cell.border=Border(bottom=side)
        for _row in sheet.iter_rows(**range_to_index('R'+ str(rowshift) +':R'+ str(rowshift+14))): #right border
           for _cell in _row:
               _cell.border=Border(right=side)
        sheet['B'+str(rowshift-3)].border=Border(left=side, top=side)
        sheet['B'+str(rowshift+14)].border=Border(left=side, bottom=side, right=side)
        sheet['R'+str(rowshift+14)].border=Border(bottom=side, right=side)
        sheet['R'+str(rowshift-1)].border=Border(right=side, bottom=side)
        sheet['R'+str(rowshift-2)].border=Border(right=side)
        sheet['R'+str(rowshift-3)].border=Border(right=side, top=side)
    
    
    elif extra_c==1:
        

        #Fill extra column:
        ##### WARNING OPTION NOT DONE YET
        
        #sheet['S'+str(rowshift-1)].value='Default'
        #sheet.column_dimensions['S'].width=6
        
        #Title format        
        sheet.merge_cells('B'+str(rowshift-3)+':R'+str(rowshift-3))
        sheet.merge_cells('D'+str(rowshift-2)+ ':R'+str(rowshift-2))
        sheet.merge_cells('B'+str(rowshift)+ ':B'+str(rowshift+14))
        sheet['B'+str(rowshift)].alignment=Alignment(horizontal='center',text_rotation=90, vertical='center')
        sheet['B'+str(rowshift-3)].alignment=Alignment(horizontal='center', vertical='center')
        sheet['D'+str(rowshift-2)].alignment=Alignment(horizontal='center', vertical='center')
        sheet['B'+str(rowshift-3)].font=Font(bold=True)
        sheet['D'+str(rowshift-2)].font=Font(italic=True)
        sheet['B'+str(rowshift)].font=Font(italic=True)
        
        #Add borders
        for _row in sheet.iter_rows(**range_to_index('B'+ str(rowshift-3) +':R'+ str(rowshift-3))): #top border
           for _cell in _row:
               _cell.border=Border(top=side)
        for _row in sheet.iter_rows(**range_to_index('B'+ str(rowshift) +':R'+ str(rowshift))): #top border 2
           for _cell in _row:
               _cell.border=Border(top=side)
        for i in ['B'+ str(rowshift-2),'B'+ str(rowshift-1)]: #left border
            a = sheet[i]
            a.border=Border(left=side)
        for _row in sheet.iter_rows(**range_to_index('B'+ str(rowshift) +':B'+ str(rowshift+14))): #left border
           for _cell in _row:
               _cell.border=Border(left=side, right=side)
        for _row in sheet.iter_rows(**range_to_index('C'+ str(rowshift+14) +':R'+ str(rowshift+14))): #bottom border
           for _cell in _row:
               _cell.border=Border(bottom=side)
        for _row in sheet.iter_rows(**range_to_index('R'+ str(rowshift) +':R'+ str(rowshift+14))): #right border
           for _cell in _row:
               _cell.border=Border(right=side)
        sheet['B'+str(rowshift-3)].border=Border(left=side, top=side)
        sheet['B'+str(rowshift+14)].border=Border(left=side, bottom=side, right=side)
        sheet['R'+str(rowshift+14)].border=Border(bottom=side, right=side)
        sheet['R'+str(rowshift-1)].border=Border(right=side, bottom=side)
        sheet['R'+str(rowshift-2)].border=Border(right=side)
        sheet['R'+str(rowshift-3)].border=Border(right=side, top=side)
        
        for _row in sheet.iter_rows(**range_to_index('R'+ str(rowshift) +':R'+ str(rowshift+14))): #Fill in white last column
           for _cell in _row:
                   fill=PatternFill(start_color='FFFFFF', end_color='FFFFFF',fill_type='solid')
                   _cell.fill=fill
    
    for j in ['C','D','E','F','G','H','I','J','K','L','M','N','O','P','Q','R']:
        sheet.column_dimensions[j].width=4    
   
   
##### ##### ##### ##### ##### ##### #####
##### #####   Main Function   ##### #####
##### ##### ##### ##### ##### ##### #####
   
   
# This function Creates a sheet to compare the two different rating
# Output: will add a sheet in the output file 
def CreateTransitionMatrix(Data, output_file_path, outlier_file_path, rowshift, typeof, name_model):
    
    Year = int(Data.snapshot.max()/100)

    ### Step 1: Select Data & Columns
    if typeof=='Prelim':
        Final_Data=SelectDataForPrelimMatrix(Data)
    elif typeof=='Migration':
        Final_Data=SelectDataForMigration(Data)
     
    #Open output
    wb =  openpyxl.load_workbook(output_file_path)
    if typeof=='Prelim':
        sheet = wb.create_sheet() 
        sheet.title = u'Prelim-Final Matrix'
    elif typeof=='Migration':
        sheet=wb[u'Prelim-Final Matrix']

        
    colshift = 4 #not changeable
    
    PDRR= range(1,16)
    
    #### Create table ### 
    
    #label the vertical axis
    for rowNum in range(15):
        row1 = rowNum+rowshift
        sheet.cell(row=row1,column=3).value=PDRR[rowNum]
    
    #label the horizontal axis
    for colNum in range(15):
        col1 = colNum+colshift
        sheet.cell(row=rowshift-1,column=col1).value=PDRR[colNum]
   
    #Add titles
    if typeof=='Prelim': 
        sheet.cell(row=rowshift-3,column=colshift-2).value='Before JBA & Final Ratings'  # modified on 20180424
        sheet.cell(row=rowshift-2,column=colshift).value='Final Rating'
        sheet.cell(row=rowshift,column=colshift-2).value='Final Rating Before JBA'
    else: 
        sheet.cell(row=rowshift-3,column=colshift-2).value='Rating Migration between '+ str(Year-1)+ ' and '+ str(Year)
        sheet.cell(row=rowshift-2,column=colshift).value=str(Year)+' Rating'
        sheet.cell(row=rowshift,column=colshift-2).value=str(Year-1)+' Rating'
    
    #Evaluate if we need extra column for default
    if typeof=='Migration':
        extra_c=1
    else:
        extra_c=0
    
    ### Fill up table ###
    aa = np.zeros((15,15))
    for rowNum in range(15):
        for colNum in range(15):
            row1=rowshift+rowNum
            col1=colshift+colNum
            temp = sum((Final_Data[u'BeforeRating']==PDRR[rowNum]) & (Final_Data[u'AfterRating']==PDRR[colNum]))
            temp=int(temp)            
            if ((not(temp==0))):        
                sheet.cell(row=row1,column=col1).value= temp
                aa[rowNum,colNum]=temp                  
    
    ###Add list of oustanding obligors
    sheet.cell(row=rowshift+4,column=colshift+len(PDRR)+2).value='Outstanding Obligors'
    sheet.cell(row=rowshift+4,column=colshift+len(PDRR)+2).font=Font(bold=True)
    df_outstanding = ComputeListOutstandingObligors_new(Final_Data)

    # remove the two lines below which were used in old outlier rule. 
    #for j in range(len(list_outstanding)):
    #    sheet.cell(row=rowshift+5+(j % 7),column=colshift+len(PDRR)+2+j//7).value=int(list_outstanding[j])
        
    ### Fill up outlier excel
    
    FillExcelOutliers(Data, outlier_file_path, df_outstanding, name_model, typeof)
    
    ### Create percentage Table ###
    sheet.cell(row=rowshift-1,column=colshift+len(PDRR)+3).value='Match'
    sheet.cell(row=rowshift-1,column=colshift+len(PDRR)+4).value='Within 1'
    sheet.cell(row=rowshift-1,column=colshift+len(PDRR)+5).value='Within 2'
    sheet.cell(row=rowshift-1,column=colshift+len(PDRR)+6).value='Outside 5'
    sheet.cell(row=rowshift-1,column=colshift+len(PDRR)+7).value='Downgrade'
    sheet.cell(row=rowshift-1,column=colshift+len(PDRR)+8).value='Upgrade'
    sheet.cell(row=rowshift-1,column=colshift+len(PDRR)+9).value='Total'
    sheet.cell(row=rowshift,column=colshift+len(PDRR)+2).value='Percentage'
    sheet.cell(row=rowshift,column=colshift+len(PDRR)+3).value=within(aa,0)/np.sum(aa)
    sheet.cell(row=rowshift,column=colshift+len(PDRR)+3).number_format='.0%'
    sheet.cell(row=rowshift,column=colshift+len(PDRR)+4).value=within(aa,1)/np.sum(aa)
    sheet.cell(row=rowshift,column=colshift+len(PDRR)+4).number_format='.0%'
    sheet.cell(row=rowshift,column=colshift+len(PDRR)+5).value=within(aa,2)/np.sum(aa)
    sheet.cell(row=rowshift,column=colshift+len(PDRR)+5).number_format='.0%'
    sheet.cell(row=rowshift,column=colshift+len(PDRR)+6).value=1-within(aa,4)/np.sum(aa)
    sheet.cell(row=rowshift,column=colshift+len(PDRR)+6).number_format='.0%'
    sheet.cell(row=rowshift,column=colshift+len(PDRR)+7).value=better_ratings(aa)/np.sum(aa)
    sheet.cell(row=rowshift,column=colshift+len(PDRR)+7).number_format='.0%'
    sheet.cell(row=rowshift,column=colshift+len(PDRR)+8).value=worse_ratings(aa)/np.sum(aa)
    sheet.cell(row=rowshift,column=colshift+len(PDRR)+8).number_format='.0%'    
    sheet.cell(row=rowshift,column=colshift+len(PDRR)+9).value=np.sum(aa)
    

    CreateStyleSheet_TransitionMatrix(sheet, rowshift, extra_c, colshift)
    
    
    wb.save(output_file_path)
    
    
    

    
        
    # -*- coding: utf-8 -*-
"""
Created on Wed Mar 16 11:57:54 2016

@author: XU79799

@author: ub71894

"""

import numpy as np
import itertools
import pandas as pd
import openpyxl
from openpyxl import Workbook
from openpyxl.styles import Border, Side, Alignment, Font
from openpyxl.formatting.rule import IconSetRule
from openpyxl.styles import PatternFill
import datetime
import pandas.tseries.offsets as pto
from openpyxl.utils.cell import coordinate_from_string, column_index_from_string


def range_to_index(range_in_str): 

    delimiter = range_in_str.find(':')
    min_cell = range_in_str[:delimiter]
    max_cell = range_in_str[(delimiter+1):]

    xy = coordinate_from_string(min_cell) 
    min_col = column_index_from_string(xy[0]) 
    min_row = xy[1]
    xy = coordinate_from_string(max_cell) 
    max_col = column_index_from_string(xy[0]) 
    max_row = xy[1]

    return({'min_col':min_col, 'max_col':max_col, 'min_row':min_row, 'max_row':max_row})


def SelectData(df): # add on 20161108
    if 'CREPROPERTYID' in df.columns.tolist():
        cols = [ u'CUSTOMERID', 'CREPROPERTYID',u'Prelim_PD', u'Final_PD']
        to_dropdup = ['CUSTOMERID', 'CREPROPERTYID', 'archive_date']
    else:
        cols = [ u'CUSTOMERID', u'Prelim_PD', u'Final_PD']
        to_dropdup = ['CUSTOMERID', 'archive_date']
        
    df.sort_values(by='archive_date', ascending=True, inplace=True)
    df.drop_duplicates(subset=to_dropdup, keep='last', inplace=True)
    df.dropna(subset=['Final_PD_Risk_Rating'], inplace=True)    

    recentsnapshot = pd.to_datetime(df.snapshot.max(), format='%Y%m') + pto.MonthEnd() # make it to the end of snapshot month
    df=df[ ((recentsnapshot- pto.DateOffset(years=1))< df[u'archive_date']) & \
    (df[u'archive_date']<= recentsnapshot)]

    df=df[cols] 
    
    return df

def CreateStyleSheet_AdditionalInfo(sheet):
    sheet.merge_cells('B2:E2')
    
    side=Side(style='thin', color="FF000000")
    for _row in sheet.iter_rows(**range_to_index('B2:E4')):
        for _cell in _row:
            fill=PatternFill(start_color='FFFFFF', end_color='FFFFFF',fill_type='solid')
            _cell.fill=fill
            _cell.border=Border(right=side, bottom=side, left=side, top=side)
            _cell.font=Font(size=9)
            _cell.alignment=Alignment(horizontal='center', vertical='center')
    
    for _row in sheet.iter_rows(**range_to_index('B2:E3')):
        for _cell in _row:
            _cell.font=Font(size=9, bold=True)
    
    sheet.column_dimensions['C'].width=13


def AdditionalInfo(Data, output_file_path, CentralTendency, numofdef=0):
    
    #Open output
    wb =  openpyxl.load_workbook(output_file_path)
    sheet = wb.create_sheet() 
    sheet.title = u'Add Info' 
    
    sheet.cell(row=2,column=2).value= 'Default Rate'
    sheet.cell(row=3,column=2).value='Num of Def'
    sheet.cell(row=3,column=3).value='Central Tendency'
    sheet.cell(row=3,column=4).value='Prelim Rating'
    sheet.cell(row=3,column=5).value='Final Rating'
    
    
    last_rating=SelectData(Data)

    
    sheet.cell(row=4,column=2).value=numofdef ##=modifiede by Yang
    sheet.cell(row=4,column=2).number_format='0'
    if 'nan'==str(CentralTendency):
        sheet.cell(row=4,column=3).value= 'NA'
    else:
        sheet.cell(row=4,column=3).value=CentralTendency
        sheet.cell(row=4,column=3).number_format='.0%'

    sheet.cell(row=4,column=4).value=(last_rating[u'Prelim_PD'].mean()/100)
    sheet.cell(row=4,column=4).number_format='.0%'
    sheet.cell(row=4,column=5).value=(last_rating[u'Final_PD'].mean()/100)
    sheet.cell(row=4,column=5).number_format='.0%'    
    
    CreateStyleSheet_AdditionalInfo(sheet)
    
    wb.save(output_file_path)
    # -*- coding: utf-8 -*-
"""

Created on Wed Mar  2 12:29:17 2016

@author: xu79799 
@author: ub71894, All work after Jun,2016
@author: ub63105, modified CreateSlideForModel Nov 3, 2016
"""


from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN


import os
import six
import copy
import pandas as pd
import datetime
import pandas.tseries.offsets as pto


def ProcessData(df): # modified on 20160830 20161028
    if 'CREPROPERTYID' in df.columns.tolist():
        cols = [ u'CUSTOMERID', 'CREPROPERTYID', u'NET_OUTSTANDING_AMOUNT', \
        u'NET_COMMITMENT_TOT', u'archive_date', u'snapshot',u'Final_PD_Risk_Rating']
        to_dropdup = ['CUSTOMERID', 'CREPROPERTYID', 'archive_date']
    else:
        cols = [ u'CUSTOMERID', u'NET_OUTSTANDING_AMOUNT', u'NET_COMMITMENT_TOT',\
        u'archive_date', u'snapshot',u'Final_PD_Risk_Rating']
        to_dropdup = ['CUSTOMERID', 'archive_date']
    #Pre-process of the data
    df=df[cols]   
    df.sort_values(by='archive_date', ascending=True, inplace=True)
    df.drop_duplicates(subset=to_dropdup, keep='last', inplace=True)
    df.dropna(subset=['Final_PD_Risk_Rating'], inplace=True)
    #Add quarter for each dates
    year_list=list(map(lambda x: x.year, list(df['archive_date'])))
    quarter_list=list(map(lambda x: ((x.month-1)//3+1), list(df['archive_date'])))
    df.loc[:,'Quarter']=list(map(lambda x,y: ('Q'+str(x)+'_'+str(y)), quarter_list,year_list))

    return df

        
def AnalyzeData(Data): # debugged on 20160830, 20161028
    
    Data=ProcessData(Data)
    recentsnapshot = pd.to_datetime(Data.snapshot.max(), format='%Y%m') + pto.MonthEnd() # make it to the end of snapshot month
    
    number_obligors=len(Data.query('snapshot==@Data.snapshot.max()'))
    numbers_obligors_rerated=len(Data[((recentsnapshot- pto.DateOffset(years=1))< Data[u'archive_date'])\
        &(Data[u'archive_date']<= recentsnapshot)])
    Commitments_tot=sum(Data.query('snapshot==@Data.snapshot.max()')[u'NET_COMMITMENT_TOT'])
    Commitments_out=sum(Data.query('snapshot==@Data.snapshot.max()')[u'NET_OUTSTANDING_AMOUNT'])

    return (number_obligors,numbers_obligors_rerated,Commitments_tot,Commitments_out)

######  ######  ######  ######  ###### 
######      Main function       ###### 
######  ######  ######  ######  ###### 

def CreateSlideForModel(Data, k, name_model):
    
    #Open Original Template
    if k==0:
        prs = Presentation( r'AuxiliaryFiles\OriginalTemplate_dup.pptx')
    else:
        prs = Presentation( r'Output_files.pptx')
     
    Slide=prs.slides[k+1]
    
   
    
    ### Documentation Table
#    t=Slide.shapes[1].table 
#    p=t.cell(1,0).text_frame.add_paragraph()
#    p.text='Python can add some predefined text_Documentation'
#    p.font.size=Pt(9)
#    
#    #### Dataset Table
#  
#    t=Slide.shapes[2].table
#    p=t.cell(1,0).text_frame.add_paragraph()
#    p.text='Python can add some predefined text_Dataset'
#    p.font.size=Pt(9)
    
    #### Title Table
   
    Slide.shapes[4].text=name_model
    
    ### Summary Table
   
    t=Slide.shapes[3].table 
    number_obligors,numbers_obligors_rerated,Commitments_tot,Commitments_out=AnalyzeData(Data)
    p=t.cell(0,1).text_frame.paragraphs[0]
    p.text='$'+str(int(Commitments_tot/1000000))
    p.font.size=Pt(9)
    p.alignment = PP_ALIGN.CENTER
    
    p=t.cell(1,1).text_frame.paragraphs[0]
    p.text='$'+str(int(Commitments_out/1000000))
    p.font.size=Pt(9)
    p.alignment = PP_ALIGN.CENTER
        
    p=t.cell(2,1).text_frame.paragraphs[0]
    p.text=str(number_obligors)
    p.font.size=Pt(9)
    p.alignment = PP_ALIGN.CENTER

    p=t.cell(3,1).text_frame.paragraphs[0]
    p.text=str(numbers_obligors_rerated)
    p.font.size=Pt(9)
    p.alignment = PP_ALIGN.CENTER

    txBox = Slide.shapes.add_textbox(Inches(0.4), Inches(6.65), Inches(1), Inches(1))
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = r"* Accuracy rate is represented by Somers'D"
    p.font.bold = True
    p.font.size = Pt(10)

    prs.save( r'Output_files.pptx')
# -*- coding: utf-8 -*-
"""

Created on Wed Mar  2 12:29:17 2016

@author: xu79799 
@author: ub71894, All work after Jun,2016

"""
import numpy as np
import itertools
import pandas as pd
import openpyxl
from openpyxl import Workbook
from openpyxl.styles import Border, Side, Alignment, Font
from openpyxl.formatting.rule import IconSetRule
from openpyxl.styles import PatternFill
import datetime
import pandas.tseries.offsets as pto
from openpyxl.utils.cell import coordinate_from_string, column_index_from_string


def range_to_index(range_in_str): 

    delimiter = range_in_str.find(':')
    min_cell = range_in_str[:delimiter]
    max_cell = range_in_str[(delimiter+1):]

    xy = coordinate_from_string(min_cell) 
    min_col = column_index_from_string(xy[0]) 
    min_row = xy[1]
    xy = coordinate_from_string(max_cell) 
    max_col = column_index_from_string(xy[0]) 
    max_row = xy[1]

    return({'min_col':min_col, 'max_col':max_col, 'min_row':min_row, 'max_row':max_row})



def Area(x, y):
    dx = np.diff(np.array(x))
    y_avg = np.array(y)
    y_avg = 0.5*(y_avg[:-1]+y_avg[1:])
    A = sum(y_avg*dx)
    return A

def AR(factor, response, sign):
    
    # convert to numpy array
    x = np.array(factor)
    y = np.array(response) #y needs to be 0 and 1
    
    # remove NaNs in either vector
    mask = ~(np.isnan(x) | np.isnan(y))
    x = x[mask]
    y = y[mask]
    
    # sort the credit risk indicator in terms of the factor values
    index = x.argsort()
    y = y[index]/sum(y)
    if sign==-1:
        y = y[::-1]    #Reverse table of y

    y_perfect = np.array(sorted(y, reverse=True)).copy() #decrease order
    
    # add in point at the origin for integration
    y_perfect = np.insert(y_perfect, 0, 0)
    y = np.insert(y, 0, 0)
    
    N = len(y)
    
    # create population % variable (x), and random and perfect model curves
    u = np.linspace(0, 1, num=N, endpoint=True)
    y_random = np.linspace(0, sum(y), num=N, endpoint=True)
    
    y = y.cumsum()    
    y_perfect = y_perfect.cumsum()
    
    # remove duplicate values and treat 
    z = pd.DataFrame(data=[np.insert(x, 0, np.nan), u, y, y_random, y_perfect], index=['x', 'u', 'y', 'y_random', 'y_perfect']).T
    z.drop_duplicates(subset=['x'], keep='last', inplace=True)

    AreaRandom = Area(z.u, z.y_random)
    AreaPerfect = Area(z.u, z.y_perfect)
    AreaActual = Area(z.u, z.y)
    
    accuracy_ratio = (AreaActual - AreaRandom)/(AreaPerfect - AreaRandom)
    
    return accuracy_ratio

def somersD(x, y, sign=1, unit_concordance=True): # debugged on 20160817
    n = len(x)
    C = 0
    D = 0
    if unit_concordance==True:
        for i, j in itertools.combinations(np.arange(n), 2):
            dx = x[i] - x[j]
            dy = y[i] - y[j]
            if (((dx<0) & (dy<0)) | ((dx>0) & (dy>0))):
                cij = 1
                dij = 1
            elif (((dx<0) & (dy>0)) | ((dx>0) & (dy<0))):
                cij = -1
                dij = 1
            else:
                cij = 0
                if dy==0: 
                    dij = 0
                else: 
                    dij=1
            C = C + cij
            D = D + dij
    else:
        for i, j in itertools.combinations(np.arange(n), 2):
            dx = x[i] - x[j]
            dy = y[i] - y[j]
            if (((dx<0) & (dy<0)) | ((dx>0) & (dy>0))):
                cij = abs(dy)
                dij = abs(dy)
            elif (((dx<0) & (dy>0)) | ((dx>0) & (dy<0))):
                cij = -abs(dy)
                dij = abs(dy)
            else:
                cij = 0
                if dy==0: 
                    dij = 0
                else: 
                    dij=1
            C = C + cij
            D = D + dij    
    
    return sign*C/D


def CreateStyleSheet_DiscriminationPower(sheet):
    
    #Work on style
    for _row in sheet.iter_rows(**range_to_index('B4:E7')):
        for _cell in _row:
            _cell.alignment=Alignment(horizontal='center',vertical='center')
            _cell.font=Font(size=9)
            fill=PatternFill(start_color='FFFFFF', end_color='FFFFFF',fill_type='solid')
            _cell.fill=fill 
    for _row in sheet.iter_rows(**range_to_index('B5:B7')):
        for _cell in _row:
            _cell.alignment=Alignment(vertical='center')

    sheet.column_dimensions['B'].width=20
    sheet.column_dimensions['C'].width=13
    sheet.column_dimensions['D'].width=13
    sheet.column_dimensions['E'].width=13
    sheet.cell(row=4,column=2).font=Font(size=9, bold=True)
    sheet.cell(row=6,column=2).font=Font(size=9, bold=True)
    sheet.cell(row=7,column=2).font=Font(size=9, bold=True)
    sheet.cell(row=5,column=3).font=Font(size=9, bold=True)
    sheet.cell(row=5,column=4).font=Font(size=9, bold=True)
    sheet.cell(row=5,column=5).font=Font(size=9, bold=True)
    
    #Add borders
    side=Side(style='thin', color="FF000000")
    sheet['B5'].border=Border(left=side)
    sheet['B6'].border=Border(left=side)
    sheet['B7'].border=Border(left=side, bottom=side)
    sheet['C7'].border=Border(bottom=side)
    sheet['D7'].border=Border(bottom=side)
    sheet['E7'].border=Border(bottom=side, right=side)
    sheet['E6'].border=Border(right=side)
    sheet['E5'].border=Border(right=side)             
    sheet['B4'].border=Border(bottom=side, left=side, top=side)
    sheet['C4'].border=Border(bottom=side, top=side)  
    sheet['D4'].border=Border(bottom=side, top=side)
    sheet['E4'].border=Border(bottom=side, right=side, top=side)    
    

# modified on 20160830. To make the calculation is based on the most recent 1 year data
def DiscriminationPower(Data, Type, output_file_path, model_type, dvlpmt_AR, dvlpmt_quantAR, dvlpmt_qualAR, relationship):
    if Type=='AR': 
        raise 'Sorry this option has not been implemented yet'
    if Type!='SomersD':
        raise 'Sorry we do not understand the type'
    
    #Open output
    wb =  openpyxl.load_workbook(output_file_path)
    sheet = wb.create_sheet() 
    sheet.title = u'Discrimination Power'
    
    #Select Data, modified on 20160830, 20161028
    if 'CREPROPERTYID' in Data.columns.tolist():
        to_dropdup = ['CUSTOMERID', 'CREPROPERTYID', 'archive_date']
    else:
        to_dropdup = ['CUSTOMERID', 'archive_date']
    
    Data.sort_values(by='archive_date', ascending=True, inplace=True)
    Data.drop_duplicates(subset=to_dropdup, keep='last', inplace=True)
    Data.dropna(subset=['Final_PD_Risk_Rating'], inplace=True)

    Year = int(Data.snapshot.max()/100)
    recentsnapshot = pd.to_datetime(Data.snapshot.max(), format='%Y%m') + pto.MonthEnd() # make it to the end of snapshot month
    Data=Data[ ((recentsnapshot- pto.DateOffset(years=1))< Data[u'archive_date']) \
    & (Data[u'archive_date']<= recentsnapshot)]
 
    #Create columns and rows:
    sheet.cell(row=4,column=2).value='Discrimination Power* - Year ' +str(Year)  
    sheet.merge_cells('B4:E4')
    sheet.cell(row=6,column=2).value='Current Portfolio'
    sheet.cell(row=7,column=2).value='Development Portfolio'
    sheet.cell(row=5,column=3).value='Overall AR'
    sheet.cell(row=5,column=4).value='Quantitative AR'
    sheet.cell(row=5,column=5).value='Qualitative AR'
        
    ### Fill Table
    if model_type=='All':
        sheet.cell(row=6,column=4).value=somersD(relationship*np.array(Data[u'Quantitative_Score']), np.array(Data[u'Final_PD']))
        sheet.cell(row=6,column=4).number_format='.0%'
        sheet.cell(row=6,column=5).value=somersD(relationship*np.array(Data[u'Qualitative_Score']), np.array(Data[u'Final_PD']))
        sheet.cell(row=6,column=5).number_format='.0%'
    if model_type=='Quantitative': #Use total score column since the quantitative one might be empty
        sheet.cell(row=6,column=4).value=somersD(relationship*np.array(Data[u'Total_Score']), np.array(Data[u'Final_PD']))
        sheet.cell(row=6,column=4).number_format='.0%'
        sheet.cell(row=6,column=5).value='NA'
    if model_type=='Qualitative': #Use total score column since the qualitative one might be empty
        sheet.cell(row=6,column=5).value=somersD(relationship*np.array(Data[u'Total_Score']), np.array(Data[u'Final_PD']))
        sheet.cell(row=6,column=5).number_format='.0%'
        sheet.cell(row=6,column=4).value='NA'

    sheet.cell(row=6,column=3).value=somersD(relationship*np.array(Data[u'Total_Score']), np.array(Data[u'Final_PD']))
    sheet.cell(row=6,column=3).number_format='.0%'
    
    #Add development information
    if 'nan'==str(dvlpmt_AR):
        sheet.cell(row=7,column=3).value='NA'
    else:
        sheet.cell(row=7,column=3).value=dvlpmt_AR
        sheet.cell(row=7,column=3).number_format='.0%'
    if 'nan'==str(dvlpmt_quantAR):
        sheet.cell(row=7,column=4).value='NA'
    else:
        sheet.cell(row=7,column=4).value=dvlpmt_quantAR
        sheet.cell(row=7,column=4).number_format='.0%'
    if 'nan'==str(dvlpmt_qualAR):
        sheet.cell(row=7,column=5).value='NA'
    else:
        sheet.cell(row=7,column=5).value=dvlpmt_qualAR
        sheet.cell(row=7,column=5).number_format='.0%'
    
    
    #Add style of the sheet
    CreateStyleSheet_DiscriminationPower(sheet)
    
    wb.save(output_file_path)
    

                
            
            
            
        
        

# -*- coding: utf-8 -*-
"""
Created on Mon Mar  7 11:43:35 2016

@author: XU79799
"""
import os
import win32com.client

def CopyOutput(output_file_path, dir_name, input_sheet,  input_selection,
               output_ppt_path, Width, Left, Top, Height, number_slide):

    excel = win32com.client.Dispatch("Excel.Application")
    wb = excel.Workbooks.Open(output_file_path)
    xl_range = wb.Sheets(input_sheet).Range(input_selection)
    xl_range.CopyPicture()
    

    excel.ActiveWorkbook.Sheets.Add( After=excel.ActiveWorkbook.Sheets(4)).Name="image_sheet"
    cht = excel.ActiveSheet.ChartObjects().Add(0,0, xl_range.Width, xl_range.Height)
    cht.Chart.Paste()
    # Export the sheet with the chart to a new file
    cht.Chart.Export(os.path.join(dir_name, 'output_fig',r'temp.bmp'),'bmp')
    excel.ActiveWorkbook.Close(False)
    excel.Quit()


    
    ##Paste into ppt
    Application = win32com.client.Dispatch("PowerPoint.Application")
    Presentation = Application.Presentations.Open(output_ppt_path)
    Slide = Presentation.Slides(number_slide)
    
    #s = Slide.Shapes.PasteSpecial()#DisplayAsIcon=0s
    #s.Width=Width
    #s.Left=Left
    #s.Top=Top
    # Save and Close the book
    Pict1 = Slide.Shapes.AddPicture(FileName=os.path.join(dir_name, 'output_fig',r'temp.bmp'),\
     LinkToFile=False, SaveWithDocument=True, Left=Left, Top=Top, Width=Width, Height=Height)
    Presentation.Save()
    Application.Quit()
    # Save and Close the book


def ExportResultsSummaryTable(output_file_path, dir_name, n_slide):
    CopyOutput(output_file_path, dir_name, "Prelim-Final Matrix",  "B7:R24", 
               os.path.join(dir_name, r'Output_files.pptx'),
               Width=293.04, Left=391.68, Top=79.92, Height=181, number_slide=n_slide)
               
    CopyOutput(output_file_path, dir_name, "Prelim-Final Matrix",  "B27:R44", 
                os.path.join(dir_name, r'Output_files.pptx'),
               Width=293.04, Left=391.68, Top=293.76, Height=181, number_slide=n_slide)
               
    CopyOutput(output_file_path,dir_name, "Prelim-Final Matrix",  "U9:AB10", 
                os.path.join(dir_name, r'Output_files.pptx'),
                Width=316.08, Left=391.68, Top=263.52, Height=21,number_slide=n_slide)
                
    CopyOutput(output_file_path, dir_name, "Prelim-Final Matrix",  "U29:AB30", 
                os.path.join(dir_name, r'Output_files.pptx'),
                Width=316.08, Left=391.68, Top=479.52, Height=21,number_slide=n_slide)

    CopyOutput(output_file_path, dir_name, "Discrimination Power",  "B4:E7", 
               os.path.join(dir_name, r'Output_files.pptx'),
               Width=310.32, Left=37.44, Top=257.2, Height=56,number_slide=n_slide)   

    CopyOutput(output_file_path, dir_name, "RLA-Override",  "B1:G7", 
               os.path.join(dir_name, r'Output_files.pptx'),
               Width=310.32, Left=37.44, Top=320, Height=120, number_slide=n_slide)

    CopyOutput(output_file_path, dir_name, "Add Info",  "B2:E4", 
               os.path.join(dir_name, r'Output_files.pptx'),
               Width=270, Left=37.44, Top=450, Height=45,number_slide=n_slide)              
               
#    CopyOutput(output_file_path, "Summary Table",  "C2:D9", 
#               os.path.join(r'C:\Users\xu79799\Documents\Python Scripts\Prototype3', r'Output_files.pptx'),
#               Width=144.72, Left=37.44, Top=95.76, number_slide=2)
#               
#    CopyOutput(output_file_path, "Summary Table",  "F12:F13", 
#               os.path.join(r'C:\Users\xu79799\Documents\Python Scripts\Prototype3', r'Output_files.pptx'),
#               Width=158.4, Left=190.08, Top=102.24, number_slide=2)
#    
#    CopyOutput(output_file_path, "Summary Table",  "F16:F17", 
#               os.path.join(r'C:\Users\xu79799\Documents\Python Scripts\Prototype3', r'Output_files.pptx'),
#               Width=158.4, Left=190.08, Top=198, number_slide=2)

# -*- coding: utf-8 -*-
"""
Created on Thu Nov  3 13:02:51 2016

@author: ub71894 (4e8e6d0b), CSG
"""
import win32com.client
import os

def PreparePPTTemp(Input_Files):

    cwd = os.getcwd()
    ppt = r'AuxiliaryFiles\OriginalTemplate.pptx'
    ppt_dup = r'AuxiliaryFiles\OriginalTemplate_dup.pptx'
    PowerPoint = win32com.client.DispatchEx( 'Powerpoint.Application' )
    Template = PowerPoint.Presentations.Open(os.path.join(cwd,ppt), False, False, False )    
    k = len(Input_Files)
    
    if k>1:
        for i in range(k-1):
            Template.Slides(2).Duplicate()
    else:
        pass
    Template.SaveAs(os.path.join(cwd,ppt_dup))
    Template.Close() 
    # -*- coding: utf-8 -*-
"""

Created on Wed Mar  2 12:29:17 2016

@author: xu79799 
@author: ub71894, All work after Jun,2016

"""

import numpy as np
import pandas as pd
import openpyxl
from openpyxl import Workbook
from openpyxl.styles import Border, Side, Alignment, Font
from openpyxl.formatting.rule import IconSetRule
from openpyxl.styles import PatternFill
import pandas.tseries.offsets as pto
from openpyxl.utils.cell import coordinate_from_string, column_index_from_string


def range_to_index(range_in_str): 

    delimiter = range_in_str.find(':')
    min_cell = range_in_str[:delimiter]
    max_cell = range_in_str[(delimiter+1):]

    xy = coordinate_from_string(min_cell) 
    min_col = column_index_from_string(xy[0]) 
    min_row = xy[1]
    xy = coordinate_from_string(max_cell) 
    max_col = column_index_from_string(xy[0]) 
    max_row = xy[1]

    return({'min_col':min_col, 'max_col':max_col, 'min_row':min_row, 'max_row':max_row})


def Color(Percentage, count, typeof): # modified on 08/17/2020 to add 'count' into coloring rule
    if Percentage==None: return ('FFFFFF')

    if (typeof=='RLA') & (count>100):
        borders=np.array([0.3,0.45])
        s=sum(Percentage>borders)
        return(['6CFA00','FFFF00','FB3030'][s])
    elif (typeof=='Override') & (count>100):
        borders=np.array([0.2,0.3])
        s=sum(Percentage>borders)
        return(['6CFA00','FFFF00','FB3030'][s])

    elif (typeof=='RLA') & (count<=100):
        s=np.where(Percentage>0.3,1,0)
        return(['6CFA00','FFFF00'][s])

    elif (typeof=='Override') & (count<=100):
        s=np.where(Percentage>0.2,1,0)
        return(['6CFA00','FFFF00'][s])


def CountRLAOverrideYear(df, Year): # modified on 20160830, 20200818
    
    #Select Data, modified on 20161108
    if 'CREPROPERTYID' in df.columns.tolist():
        to_dropdup = ['CUSTOMERID', 'CREPROPERTYID', 'archive_date']
    else:
        to_dropdup = ['CUSTOMERID', 'archive_date']
    
    df.sort_values(by='archive_date', ascending=True, inplace=True)
    df.drop_duplicates(subset=to_dropdup, keep='last', inplace=True)
    df.dropna(subset=['Final_PD_Risk_Rating'], inplace=True)


    diff = int(df.snapshot.max()/100) - Year
    recentsnapshot = pd.to_datetime(df.snapshot.max(), format='%Y%m') + pto.MonthEnd() # make it to the end of snapshot month
    df=df[ ((recentsnapshot- pto.DateOffset(years=1+diff))< df[u'archive_date']) & \
    (df[u'archive_date']<= (recentsnapshot- pto.DateOffset(years=diff)))]
    
    #calculate count, modified on 20200818
    RLA_count = sum((~df[r'RLA_Notches'].isnull()) & (df[r'RLA_Notches']!=0))
    RLA=float(RLA_count)/float(df.shape[0])
    Override_count = sum(~df[r'Override_Action'].isnull() &  df[r'Override_Action']!=0)
    Override=float(Override_count)/float(df.shape[0])
    Override_plus=float(sum(df[~df[r'Override_Action'].isnull()] [ \
        r'Override_Action']>0))/float(df.shape[0])
    Override_neg=float(sum(df[~df[r'Override_Action'].isnull()] [ \
        r'Override_Action']<0))/float(df.shape[0])
    Count=df.shape[0]
    
    return (RLA, RLA_count, Override, Override_count, Override_plus, Override_neg, Count)
    

def CreateStyleSheet_RLA_Override(sheet):  # modified on 08/17/2020
    
    for col in range(3,8):
        sheet.cell(row=3,column=col).number_format='.0%'
        sheet.cell(row=4,column=col).number_format='.0%'
        sheet.cell(row=5,column=col).number_format='.0%'    
        sheet.cell(row=6,column=col).number_format='.0%'
        
    #Work on style
    for _row in sheet.iter_rows(**range_to_index('B1:G7')):
        for _cell in _row:
            _cell.alignment=Alignment(horizontal='center',vertical='center')
            _cell.font=Font(size=9)
            #fill=PatternFill(start_color='FFFFFF', end_color='FFFFFF',fill_type='solid')
            #_cell.fill=fill 
    #for _row in sheet.iter_rows(**range_to_index('C3:G3')): #Adjust color for cells
    #    for _cell in _row:
    #        color=Color(_cell.value, 'RLA')
    #        fill=PatternFill(start_color=color, end_color=color,fill_type='solid')
    #        _cell.fill=fill 
    #for _row in sheet.iter_rows(**range_to_index('C4:G4')): #Adjust color for cells
    #    for _cell in _row:
    #        color=Color(_cell.value, 'Override')
    #        fill=PatternFill(start_color=color, end_color=color,fill_type='solid')
    #        _cell.fill=fill
            
                    
    sheet.row_dimensions[2].height=17
    sheet.row_dimensions[3].height=17
    sheet.row_dimensions[4].height=17   
    sheet.row_dimensions[5].height=17    
    sheet.row_dimensions[6].height=17    
    sheet.row_dimensions[7].height=17
    sheet.column_dimensions['B'].width=14
    sheet.column_dimensions['C'].width=9
    sheet.column_dimensions['D'].width=9
    sheet.column_dimensions['E'].width=9   
    sheet.column_dimensions['F'].width=9
    sheet.column_dimensions['G'].width=9
    
    for i in ['B1', 'B2','B3','B4', 'B5','B6', 'B7']:
        a = sheet[i]
        a.font=Font(bold=True, size=9)
        a.alignment=Alignment(vertical='center', horizontal='left')
    for i in ['B2', 'C2','D2','E2', 'F2', 'G2']:
        a = sheet[i]
        a.font=Font(bold=True, size=9)
        a.alignment=Alignment(vertical='center', horizontal='center')
    sheet.merge_cells('B1:G1')
    sheet['B1'].alignment=Alignment(horizontal='center')
    
    #### Add borders ###
    side=Side(style='thin', color="FF000000")
    
    for i in ['C1','D1','E1', 'F1']: #top border
        a = sheet[i]
        a.border=Border(bottom=side, top=side)
    sheet['G1'].border=Border(bottom=side, right=side, top=side)
    sheet['B1'].border=Border(bottom=side, left=side, top=side)
    
    
    for i in ['C2','D2','E2', 'F2']: #top border
        a = sheet[i]
        a.border=Border(bottom=side, top=side)
    sheet['G2'].border=Border(bottom=side, right=side, top=side)
    sheet['B2'].border=Border(bottom=side, left=side, top=side, right=side)
    
    sheet['B3'].border=Border(left=side, top=side, right=side)
    sheet['G3'].border=Border(right=side, top=side)
    
    sheet['B4'].border=Border(left=side, right=side)
    sheet['G4'].border=Border(right=side)
    
    for i in ['C5','D5','E5', 'F5']: #top border
        a = sheet[i]
        a.border=Border(top=side)
    sheet['B5'].border=Border(top=side, left=side, right=side)
    sheet['G5'].border=Border(right=side, top=side)    
    
    for i in ['C6','D6','E6', 'F6']: #top border
        a = sheet[i]
        a.border=Border(bottom=side)
    sheet['B6'].border=Border(bottom=side, left=side, right=side)
    sheet['G6'].border=Border(bottom=side, right=side)      
    
    sheet['B7'].border=Border(bottom=side, right=side, left=side, top=side)      
    sheet['C7'].border=Border(bottom=side, left=side, top=side) 
    sheet['D7'].border=Border(bottom=side, top=side)       
    sheet['E7'].border=Border(bottom=side, top=side)  
    sheet['F7'].border=Border(bottom=side, top=side)
    sheet['G7'].border=Border(bottom=side, top=side, right=side)  


def CreateStatistics_RLA_Override(Data, output_file_path):# modified on 20160830, 20200818
        
    #Open output
    wb =  openpyxl.load_workbook(output_file_path)
    sheet = wb.create_sheet() 
    sheet.title = u'RLA-Override'
    Year = int(Data.snapshot.max()/100)

    #Create output Table
    sheet.cell(row=1,column=2).value='RLA & Override'
    sheet.cell(row=3,column=2).value='RLA'
    sheet.cell(row=4,column=2).value='Overall Override'
    sheet.cell(row=5,column=2).value='Override Positive'
    sheet.cell(row=6,column=2).value='Override Negative'
    sheet.cell(row=7,column=2).value='Count'

   
    #Fill up RLA & Override
    col=3
    tot_yr = len(Data.snapshot.unique())  # added on 03/28/2019 since some model\
                                          # has less than 5 year history                             
    for i in range(tot_yr):
        RLA,RLA_count, Override,Override_count, Override_plus, Override_neg, Count=CountRLAOverrideYear(Data, Year-i)
        sheet.cell(row=2,column=col).value=Year-i

        # change coloring rule on 20200818
        sheet.cell(row=3,column=col).value=RLA
        color=Color(RLA, RLA_count, 'RLA')
        fill=PatternFill(start_color=color, end_color=color,fill_type='solid')
        sheet.cell(row=3,column=col).fill=fill 

        sheet.cell(row=4,column=col).value=Override
        color=Color(Override, Override_count, 'Override')
        fill=PatternFill(start_color=color, end_color=color,fill_type='solid')
        sheet.cell(row=4,column=col).fill=fill

        sheet.cell(row=5,column=col).value=Override_plus   
        sheet.cell(row=6,column=col).value=Override_neg
        sheet.cell(row=7,column=col).value=Count
        col+=1
    
    #Work on Sheet format & Style
    CreateStyleSheet_RLA_Override(sheet)
    
    #del original sheet
    del wb['Sheet']    
      
    
    #Save excel file
    wb.save(output_file_path)





# -*- coding: utf-8 -*-
"""
Created on Fri Oct 14 13:46:54 2016

@author: ub71894
"""

from .Initiate import *
from .RLAOverride import *
from .CalculateMatrix import *
from .DiscriminationPower import *
from .CreateAdditionalInfo import *
from .CreatePptTable import *
from .ExportResult import *

__version__ = "0.2.0"
['CalculateMatrix.py', 'CreateAdditionalInfo.py', 'CreatePptTable.py', 'DiscriminationPower.py', 'ExportResult.py', 'Initiate.py', 'RLAOverride.py', '__init__.py']
[516, 627, 750, 1009, 1095, 1121, 1358, 1374]
